# agent/export/renderer.py
from __future__ import annotations
import re
from copy import deepcopy
from pathlib import Path
from typing import BinaryIO, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from .spec import SlideDeckSpec, SlideSpec, TableSpec

# python-pptx PP_PLACEHOLDER.SLIDE_NUMBER == 13 (sldNum). 명시 상수로 매직넘버 회피.
PH_TYPE_SLIDE_NUMBER = 13


# §13-3 v2 amendment (2026-05-08): 템플릿 재정비 후 placeholder 기반 식별로 전환.
# placeholder 식별 규칙: layout name + top 좌표 기준 (idx 단독 사용 금지 — idx 10/11/12 가
# 다른 layout 에서 DATE/FOOTER/SLIDE_NUMBER 와 충돌하므로 top 좌표가 결정 키).
EMU_PER_CM = 360000

# layout 2 (SECTION_HEADER) placeholder top 좌표 (cm)
SECTION_HEADER_NUMBER_TOP_CM = 5.0     # 챕터 번호 (예: "01")
SECTION_HEADER_TITLE_TOP_CM = 10.5     # 챕터 제목
SECTION_HEADER_SUBTITLE_TOP_CM = 13.0  # 부제 (optional)

# layout 0 (TITLE) 부제 강제 재배치 좌표 — 마스터에서 이미 정렬됐으면 no-op 안전망.
# 4개 모두 명시 필수: left/width 만 set 시 spPr 신설로 top/height master inherit 끊김 → 0 추락.
TITLE_SUBTITLE_LEFT_CM = 2.0
TITLE_SUBTITLE_TOP_CM = 9.09
TITLE_SUBTITLE_WIDTH_CM = 29.87
TITLE_SUBTITLE_HEIGHT_CM = 1.50

# layout 5 ('제목만') 표 그릴 영역 (layout 1 OBJECT 영역과 동일 좌표 — 디자인 일관성).
TABLE_AREA_LEFT_CM = 2.0
TABLE_AREA_TOP_CM = 4.0
TABLE_AREA_WIDTH_CM = 29.87
TABLE_AREA_HEIGHT_CM = 12.5

LAYOUT_TITLE_TABLE = 5  # 'TITLE_TABLE' — 표 슬라이드 dispatch 대상 (v3: 'TITLE_CONTENT' 와 제목 좌표 통일)

# §13-10 표 스타일 상수 (코드 정의 — python-pptx 가 표 스타일 갤러리 적용 미지원하기 때문).
# 마스터 테마는 폰트 종류·색 팔레트만 상속, 폰트 크기·헤더 강조·컬럼 너비는 코드 명시.
TABLE_HEADER_FONT_NAME = "Pretendard"
TABLE_HEADER_FONT_SIZE = Pt(12)
TABLE_HEADER_FONT_SIZE_FALLBACK = Pt(11)  # 행 수 ≥ TABLE_FALLBACK_THRESHOLD_ROWS
TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x1A)  # 차콜
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)  # 흰색

TABLE_BODY_FONT_NAME = "Pretendard"
TABLE_BODY_FONT_SIZE = Pt(10)
TABLE_BODY_FONT_SIZE_FALLBACK = Pt(9)  # 행 수 ≥ TABLE_FALLBACK_THRESHOLD_ROWS
TABLE_BODY_FG = RGBColor(0x1A, 0x1A, 0x1A)

# 행 높이는 명시적 고정 X (PowerPoint 자동 맞춤 활성화). 최소값만 설정.
TABLE_ROW_HEIGHT_MIN_CM = 0.6
# 컬럼 너비 휴리스틱: 셀 텍스트 평균 길이 비례 분배. 최소 너비 보장.
TABLE_COL_MIN_WIDTH_CM = 2.5
# 행 수 임계 — 이상이면 폰트 크기 fallback (헤더 11pt / 본문 9pt) 적용.
TABLE_FALLBACK_THRESHOLD_ROWS = 8


def render_deck(
    spec: SlideDeckSpec,
    *,
    template_path: Union[str, Path],
    out: Union[str, Path, BinaryIO],
) -> Optional[Path]:
    """SlideDeckSpec → .pptx 결정론적 렌더 (LLM 호출 0).

    layout_id 분기:
      - 0 (TITLE): idx=0 CENTER_TITLE = title, idx=1 SUBTITLE = body or topic_title
      - 1 (TITLE_CONTENT) + table: layout 5 ('TITLE_TABLE') 로 dispatch + add_table
      - 1 (TITLE_CONTENT) bullets/body: idx=0 TITLE, idx=1 OBJECT 채움
      - 2 (SECTION_HEADER): placeholder 3개 (top 좌표 식별) — 번호/제목/부제

    §13-12-2: out 인자 BinaryIO 지원 (R1 시그니처 확장).
      - str/Path 인자: 파일 경로. 부모 디렉터리 자동 생성. Path 반환.
      - BinaryIO (예: io.BytesIO): 메모리 스트림. python-pptx Presentation.save() native 지원.
        HTTP 응답 직접 streaming 용도. None 반환.
      이 패턴은 향후 export 형식 추가 (PDF deck 등) 시 동일하게 재사용.
    """
    template_path = Path(template_path)
    if not template_path.exists():
        raise FileNotFoundError(f"template_path not found: {template_path}")

    prs = Presentation(str(template_path))
    _clear_template_slides(prs)

    # §13-15 (2026-05-12): SECTION_HEADER chapter_number 를 LLM 산출 (title prefix) 의존에서
    # deterministic counter 로 전환. + subtitle 의 layout default inheritance 차단.
    section_header_counter = 0
    for i, s in enumerate(spec.slides):
        if s.layout_id == 1 and s.table:
            # 표 슬라이드 — layout 5 ('TITLE_TABLE') 로 dispatch (placeholder/표 좌표 겹침 회피)
            layout = prs.slide_layouts[LAYOUT_TITLE_TABLE]
            slide = prs.slides.add_slide(layout)
            _render_title_table_slide(slide, s)
        else:
            layout = prs.slide_layouts[s.layout_id]
            slide = prs.slides.add_slide(layout)
            if s.layout_id == 0:
                _render_title_slide(slide, s, default_subtitle=spec.topic_title)
            elif s.layout_id == 1:
                _render_content_slide(slide, s)
            elif s.layout_id == 2:
                section_header_counter += 1
                subtitle_fallback = _next_titled_content_summary(spec.slides, i)
                _render_section_header_slide(
                    slide, s,
                    chapter_number=section_header_counter,
                    subtitle_fallback=subtitle_fallback,
                )
        # python-pptx add_slide() 가 SLIDE_NUMBER placeholder 를 자동 상속하지 않음 →
        # layout 에 정의돼 있으면 sp XML 을 deep-copy 해서 슬라이드 spTree 에 명시 추가.
        # SECTION_HEADER (layout 2) 는 layout 자체에 SLIDE_NUMBER 가 없어 자동 skip — 의도적 제외.
        _ensure_slide_number(slide, layout)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes

    if isinstance(out, (str, Path)):
        out_path = Path(out)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return out_path
    prs.save(out)
    return None


def _clear_template_slides(prs) -> None:
    """템플릿에 미리 들어있던 starter slide 삭제 (layout 보존)."""
    sld_id_lst = prs.slides._sldIdLst
    rid_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sld_id in list(sld_id_lst):
        rid = sld_id.get(rid_attr)
        prs.part.drop_rel(rid)
        sld_id_lst.remove(sld_id)


def _placeholder_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _ensure_slide_number(slide, layout) -> None:
    """layout 에 SLIDE_NUMBER placeholder 가 있으면 슬라이드 인스턴스로 명시 복사.

    python-pptx add_slide() 동작: TITLE/CENTER_TITLE/BODY/OBJECT 등 일반 placeholder 는
    layout → 슬라이드 자동 상속하지만, SLIDE_NUMBER 만 자동 상속 제외 (검증 결과).
    layout 의 sp XML (자동 필드 <a:fld type="slidenum"> 포함) 을 deep-copy 해서
    슬라이드 spTree 에 추가하면 PowerPoint 가 슬라이드 번호 자동 표시.

    cNvPr id 충돌 회피: layout sp 의 id (보통 4 or 5) 가 슬라이드 spTree 의 다른
    shape 과 겹칠 수 있어 max(existing) + 1 로 재할당.
    """
    # layout 에서 SLIDE_NUMBER placeholder 찾기
    layout_sp = None
    for ph in layout.placeholders:
        try:
            if int(ph.placeholder_format.type) == PH_TYPE_SLIDE_NUMBER:
                layout_sp = ph._element
                break
        except Exception:
            continue
    if layout_sp is None:
        return  # SECTION_HEADER 등 SLIDE_NUMBER 미정의 layout — 의도적 제외

    # 이미 슬라이드에 SLIDE_NUMBER 가 있으면 skip (중복 방지·idempotent)
    for ph in slide.placeholders:
        try:
            if int(ph.placeholder_format.type) == PH_TYPE_SLIDE_NUMBER:
                return
        except Exception:
            continue

    sp_clone = deepcopy(layout_sp)

    # cNvPr id 재할당 — 슬라이드 spTree 내 모든 cNvPr id max + 1
    spTree = slide.shapes._spTree
    existing_ids: list[int] = []
    for cnv in spTree.iter(qn("p:cNvPr")):
        try:
            existing_ids.append(int(cnv.get("id", "0")))
        except (ValueError, TypeError):
            continue
    new_id = max(existing_ids, default=1) + 1

    cnv_clone = sp_clone.find(".//" + qn("p:cNvPr"))
    if cnv_clone is not None:
        cnv_clone.set("id", str(new_id))

    spTree.append(sp_clone)


def _placeholder_by_top_cm(slide, top_cm: float, *, tol_cm: float = 0.5):
    """top 좌표 ± tol_cm 안에 있는 placeholder 1개 반환. idx 단독 식별 금지(SECTION_HEADER 의
    BODY placeholder idx 10/11/12 가 타 layout 의 DATE/FOOTER/SLIDE_NUMBER 와 충돌)."""
    target = top_cm * EMU_PER_CM
    tol = tol_cm * EMU_PER_CM
    for ph in slide.placeholders:
        try:
            if ph.top is not None and abs(ph.top - target) <= tol:
                return ph
        except Exception:
            continue
    return None


_CHAPTER_NUM_RE = re.compile(r"^\s*(\d+)\.\s*(.+)$")


def _split_chapter_title(title: str) -> tuple[str, str]:
    """`'1. Executive Summary'` → `('01', 'Executive Summary')`. 매칭 안되면 ('', title)."""
    m = _CHAPTER_NUM_RE.match(title or "")
    if not m:
        return "", (title or "").strip()
    n_str, rest = m.group(1), m.group(2).strip()
    return f"{int(n_str):02d}", rest


def _render_title_slide(slide, s: SlideSpec, *, default_subtitle: str) -> None:
    title_ph = _placeholder_by_idx(slide, 0)
    subtitle_ph = _placeholder_by_idx(slide, 1)
    if title_ph is not None:
        title_ph.text = s.title
    if subtitle_ph is not None:
        subtitle_ph.text = s.body or default_subtitle
        # 마스터 부제가 left=4.2cm 어긋난 케이스 대비 강제 재배치 (이미 정렬됐으면 no-op 안전).
        # 4개 좌표 모두 명시 필수 — left/width 만 set 시 master inherit 끊겨 top/h=0 으로 추락.
        subtitle_ph.left = Cm(TITLE_SUBTITLE_LEFT_CM)
        subtitle_ph.top = Cm(TITLE_SUBTITLE_TOP_CM)
        subtitle_ph.width = Cm(TITLE_SUBTITLE_WIDTH_CM)
        subtitle_ph.height = Cm(TITLE_SUBTITLE_HEIGHT_CM)


def _render_content_slide(slide, s: SlideSpec) -> None:
    """layout 1 (TITLE_CONTENT) — bullets > body 만 처리. table 분기는 layout 5 dispatch."""
    title_ph = _placeholder_by_idx(slide, 0)
    object_ph = _placeholder_by_idx(slide, 1)
    if title_ph is not None:
        title_ph.text = s.title
    if object_ph is None:
        return
    if s.bullets:
        _fill_bullets(object_ph, s.bullets)
    elif s.body:
        object_ph.text = s.body


def _fill_bullets(placeholder, bullets) -> None:
    tf = placeholder.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b


def _render_title_table_slide(slide, s: SlideSpec) -> None:
    """layout 5 ('TITLE_TABLE') — 제목 placeholder 만 채우고 표는 별도 add_table 로 그림.
    layout master 가 TITLE_CONTENT 와 동일 제목 좌표 + 버건디 강조 라인 + 슬라이드 번호 보유.

    §13-10 스타일 적용:
      - 헤더 행: Pretendard Bold 12pt, RGB(255,255,255) 글자, RGB(26,26,26) 차콜 배경.
      - 본문 행: Pretendard 10pt, RGB(26,26,26) 글자.
      - 행 수 ≥ 8: 헤더 11pt / 본문 9pt fallback (placeholder 영역 12.5cm 보호).
      - 컬럼 너비: 셀 텍스트 평균 길이 비례 (최소 2.5cm 보장, 합계 29.87cm 유지).
      - 행 높이: 최소만 설정 (PowerPoint 자동 맞춤 활성화).
    """
    title_ph = _placeholder_by_idx(slide, 0)
    if title_ph is not None:
        title_ph.text = s.title
    if not s.table or not s.table.header:
        return
    n_rows = len(s.table.rows) + 1
    n_cols = len(s.table.header)
    if n_cols == 0:
        return

    slide.shapes.add_table(
        n_rows, n_cols,
        Cm(TABLE_AREA_LEFT_CM), Cm(TABLE_AREA_TOP_CM),
        Cm(TABLE_AREA_WIDTH_CM), Cm(TABLE_AREA_HEIGHT_CM),
    )
    table = slide.shapes[-1].table

    use_fallback = n_rows >= TABLE_FALLBACK_THRESHOLD_ROWS
    header_size = TABLE_HEADER_FONT_SIZE_FALLBACK if use_fallback else TABLE_HEADER_FONT_SIZE
    body_size = TABLE_BODY_FONT_SIZE_FALLBACK if use_fallback else TABLE_BODY_FONT_SIZE

    for c, h in enumerate(s.table.header):
        cell = table.cell(0, c)
        cell.text = h
        _style_cell_header(cell, header_size)

    for r, row in enumerate(s.table.rows):
        for c, cell_text in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r + 1, c)
            cell.text = cell_text
            _style_cell_body(cell, body_size)

    _set_column_widths(table, s.table, n_cols)

    for r in range(n_rows):
        table.rows[r].height = Cm(TABLE_ROW_HEIGHT_MIN_CM)


def _style_cell_header(cell, size) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = TABLE_HEADER_FONT_NAME
            run.font.bold = True
            run.font.size = size
            run.font.color.rgb = TABLE_HEADER_FG


def _style_cell_body(cell, size) -> None:
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = TABLE_BODY_FONT_NAME
            run.font.size = size
            run.font.color.rgb = TABLE_BODY_FG


def _set_column_widths(table, table_spec: TableSpec, n_cols: int) -> None:
    """컬럼 너비를 셀 텍스트 평균 길이에 비례 분배.

    알고리즘:
      1) 각 컬럼의 (header + 모든 본문 셀) 평균 char 길이 계산.
      2) 비율 × TABLE_AREA_WIDTH_CM 으로 ideal width 계산.
      3) 최소 너비 TABLE_COL_MIN_WIDTH_CM 보장.
      4) 합계가 area width 초과 시 비례 축소 (총합 = area width 유지).
    """
    avg_lengths: list[float] = []
    for c in range(n_cols):
        lengths = [len(table_spec.header[c])]
        for row in table_spec.rows:
            if c < len(row):
                lengths.append(len(row[c]))
        avg_lengths.append(sum(lengths) / len(lengths) if lengths else 1.0)

    total_len = sum(avg_lengths)
    if total_len <= 0:
        return  # all empty — leave default

    total_w = TABLE_AREA_WIDTH_CM
    ideal = [(l / total_len) * total_w for l in avg_lengths]
    adjusted = [max(w, TABLE_COL_MIN_WIDTH_CM) for w in ideal]
    # 항상 합계 = total_w 유지 (확대·축소 모두 비례 적용).
    # 짧은 텍스트만 있어 sum_adj < total_w 인 경우에도 area 를 꽉 채우도록.
    sum_adj = sum(adjusted)
    if sum_adj > 0:
        scale = total_w / sum_adj
        adjusted = [w * scale for w in adjusted]

    for c in range(n_cols):
        table.columns[c].width = Cm(adjusted[c])


def _next_titled_content_summary(slides, i: int) -> str:
    """SECTION_HEADER (slides[i]) 다음 첫 TITLE_CONTENT slide 의 첫 bullet 또는 body 한 줄.

    §13-15 (2026-05-12): subtitle_fallback 산출 — Phase 2 최소 변경 deterministic 추출.
    다음 SECTION_HEADER 만나면 stop (인접 빈 section 케이스). LLM call 0 추가.
    호출처 (render_deck) 가 결과를 _render_section_header_slide 의 subtitle_fallback 로 전달.
    """
    for j in range(i + 1, len(slides)):
        ns = slides[j]
        if ns.layout_id == 2:
            return ""
        if ns.layout_id == 1:
            if ns.bullets and ns.bullets[0]:
                return ns.bullets[0]
            if ns.body and ns.body.strip():
                return ns.body.split("\n", 1)[0]
            return ""
    return ""


def _render_section_header_slide(
    slide,
    s: SlideSpec,
    *,
    chapter_number: int,
    subtitle_fallback: str = "",
) -> None:
    """layout 2 — placeholder 3개 (top 좌표 식별):
      - top≈5.0cm  → 챕터 번호 (deterministic counter "01"~"99", §13-15 fix)
      - top≈10.5cm → 챕터 제목 (LLM 산출 title 의 prefix 정리)
      - top≈13.0cm → 부제 — 3-단계 우선순위 (s.body → subtitle_fallback → NBSP)

    §13-15 (2026-05-12) catch 33: layout master placeholder default 텍스트가 PowerPoint
    inheritance 로 노출되는 버그 차단. 빈 문자열 set 만으로는 차단 불가 — NBSP 또는 명시적
    텍스트 set 필수. chapter_number 는 LLM 산출 (title prefix) 의존성 제거하고 deterministic.
    """
    number_str = f"{chapter_number:02d}"
    _, clean_title = _split_chapter_title(s.title)

    num_ph = _placeholder_by_top_cm(slide, SECTION_HEADER_NUMBER_TOP_CM)
    title_ph = _placeholder_by_top_cm(slide, SECTION_HEADER_TITLE_TOP_CM)
    subtitle_ph = _placeholder_by_top_cm(slide, SECTION_HEADER_SUBTITLE_TOP_CM)

    if num_ph is not None:
        num_ph.text_frame.text = number_str
    if title_ph is not None:
        title_ph.text_frame.text = clean_title or s.title
    if subtitle_ph is not None:
        # subtitle 우선순위 (§13-15 명시화):
        #   1. s.body — planner 명시적 산출 (None 아니고 strip 후 비어있지 않음)
        #   2. subtitle_fallback — renderer 자동 추출 (다음 TITLE_CONTENT 의 첫 bullet/body)
        #   3. NBSP — edge case, layout default inheritance 차단만
        if s.body is not None and s.body.strip():
            subtitle_text = s.body.strip()
        elif subtitle_fallback.strip():
            subtitle_text = subtitle_fallback.strip()
        else:
            subtitle_text = " "
        subtitle_ph.text_frame.text = subtitle_text


__all__ = ["render_deck"]
