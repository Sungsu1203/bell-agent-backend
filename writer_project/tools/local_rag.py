# tools/local_rag.py
from __future__ import annotations
import logging
logger = logging.getLogger(__name__)

import os, re, json, glob, hashlib
from pathlib import Path
import threading
from datetime import datetime
from typing import Any, List, Tuple, Optional, Callable
from urllib.parse import unquote, quote, urlencode
from typing import Dict
from fnmatch import fnmatch
from collections import Counter

from langchain_core.documents import Document


# -----------------------------------------------------------------------------
# SSoT alignment (web_rag/utils):
# - metadata["source"]        : canonical file:// URI (with fragment for part/index/chunk)
# - metadata["source_version"]: mtime string (used by web_rag/utils.make_doc_id in ingest_vector)
# - DO NOT generate local-only doc_id based on custom __v_/__c_ encoding
# -----------------------------------------------------------------------------
from tools.web_rag.utils import _normalize_canonical_url


def _build_local_source(p: Path, *, part: str | None = None, index: int | None = None, chunk: int | None = None) -> str:
    """
    Build canonical 'source' for local chunks:
      file:///.../file.pdf#part=Slide&index=3&chunk=12
    Fragment is preserved by utils._normalize_canonical_url (file:// policy).
    """
    # 1) raw file URI 생성 (OS resolve 기반으로 최대한 단일화)
    base = _file_uri(str(p))

    # 2) fragment는 "항상 동일한 순서/인코딩"으로 구성한다.
    #    (file:// 정책상 utils._normalize_canonical_url()이 fragment를 그대로 보존하므로,
    #     표준화는 여기서 강제해야 doc_id가 흔들리지 않는다.)
    frag: dict[str, str] = {}
    if part is not None:
        frag["part"] = str(part)
    if index is not None:
        frag["index"] = str(index)
    if chunk is not None:
        frag["chunk"] = str(chunk)

    if frag:
        # quote keys/values via urlencode (stable encoding)
        base = base + "#" + urlencode(frag, doseq=False, quote_via=quote)

    # IMPORTANT:
    # Do NOT canonicalize here.
    # Canonicalization is SSoT in tools.web_rag.utils.make_doc_id().
    return base

# ── explicit module export list (typed) ─────────────────────────────────────
# 현재 파일에서 외부로 제공하는 퍼블릭 API만 나열합니다.
__all__: list[str] = [
    "build_webjson_from_local",
    "ingest_local_files",
    "add_local_findings_to_chroma",
    "quick_ingest_findings",
]


from core.config import CFG
from core.config import reload_config as reload_config  # 런타임 갱신 허용
from core.paths import (
    current_path,
    research_base_dir,
    research_topic_dir,
    research_resources_dir,   # ← 표준 리소스 경로 헬퍼 사용
)

# web_rag 유틸: 네임스페이스/디렉터리 규칙과 일치시킴
# - Optional로 선선언 후, 실제 구현은 별도 alias를 통해 대입하여 재정의(no-redef) 회피
_wr_resolve_persist_dir: Optional[Callable[[str, Optional[str]], str]] = None
_wr_sanitize_ns: Optional[Callable[[str], str]] = None
try:
    from tools.web_rag.utils import _resolve_persist_dir as _wr_resolve_persist_dir_impl
    from tools.web_rag.utils import sanitize_ns as _wr_sanitize_ns_impl
    _wr_resolve_persist_dir = _wr_resolve_persist_dir_impl
    _wr_sanitize_ns = _wr_sanitize_ns_impl
except Exception:
    # 안전 폴백: 그대로 None 유지
    pass

# ──────────────────────────────────────────────────────────────────────────────
# Optional dependencies: 클래스/함수 핸들을 Any로 보관(없으면 None)
# ──────────────────────────────────────────────────────────────────────────────
try:
    from bs4 import BeautifulSoup as _BeautifulSoup
    BeautifulSoup: Optional[Any] = _BeautifulSoup
except Exception:
    BeautifulSoup = None
    logger.debug("BeautifulSoup not available; falling back to regex for HTML parsing.")

try:
    # 프로젝트 표준: PyPDF2
    from PyPDF2 import PdfReader as _PdfReader
    PdfReader: Optional[Any] = _PdfReader
except Exception:
    PdfReader = None
    logger.debug("PyPDF2 not available; PDF extraction disabled.")

try:
    import docx as _docx  # python-docx
    docx: Optional[Any] = _docx
except Exception:
    docx = None
    logger.debug("python-docx not available; .docx extraction disabled.")

# (선택) Unstructured: PPTX/XLSX 등 포맷을 자동 분해
# unstructured: 있으면 함수/클래스 핸들을 보관, 없으면 None
try:
    from unstructured.partition.auto import partition as _unstructured_partition  # noqa: F401
    partition: Optional[Any] = _unstructured_partition
except Exception:
    partition = None
    logger.debug("unstructured.partition.auto not available.")
try:
    from unstructured.documents.elements import Table as _UnstructuredTable  # noqa: F401
    Table: Optional[Any] = _UnstructuredTable
except Exception:
    Table = None
    logger.debug("unstructured.documents.elements.Table not available.")

# (선택) openpyxl: XLSX → TSV 직렬화
# - 모듈/None 겸용을 위해 Optional[Any]로 선언(모듈 타입 고정 금지)
_OpenPyXL: Optional[Any]
try:
    import openpyxl as _openpyxl  # noqa: F401
    _OpenPyXL = _openpyxl
except Exception:
    _OpenPyXL = None
    logger.debug("openpyxl not available; XLSX TSV extraction will fallback to unstructured.")

# mypy 경고 방지용: openpyxl stubs 미설치 시 무시
# (타입 전용 import guard)
if False:  # pragma: no cover
    import types_openpyxl  # type: ignore[import-not-found]

# (선택) python-pptx: 슬라이드 단위 샘플 추출용
# python-pptx: 클래스 핸들을 Optional[Any]로 유지
try:
    from pptx import Presentation as _pptx_Presentation  # noqa: F401
    _PptxPresentation: Optional[Any] = _pptx_Presentation
except Exception:
    _PptxPresentation = None
    logger.debug("python-pptx not available; PPTX sampling disabled.")

# ──────────────────────────────────────────────────────────────────────────────
# reload_config() 1회 호출 가드 (P1-5)
#  - 기본: 최초 1회만 호출
#  - 디버깅/테스트 시 LOCAL_RAG_RELOAD_EACH_CALL=1 로 매 호출 허용
_RELOAD_ONCE_FLAG = False
_RELOAD_LOCK = threading.Lock()

# ──────────────────────────────────────────────────────────────────────────────
# 내부 유틸 (CFG 우선 → ENV 폴백 단일 진입)
# ──────────────────────────────────────────────────────────────────────────────
# ※ 다른 모듈과의 중복 정의를 피하기 위해, 이미 전역에 있으면 재정의하지 않음
if "_cfg_str" not in globals():
    def _cfg_str(name: str, default: str = "") -> str:
        try:
            v = getattr(CFG, name)
            if v is None:
                raise AttributeError
            s = str(v).strip()
            return s if s != "" else default
        except Exception:
            s = (os.getenv(name, "") or "").strip()
            return s if s != "" else default

if "_cfg_int" not in globals():
    def _cfg_int(name: str, default: int = 0) -> int:
        try:
            v = getattr(CFG, name)
            if v is None or str(v).strip() == "":
                raise AttributeError
            return int(str(v).strip())
        except Exception:
            ev = (os.getenv(name, "") or "").strip()
            try:
                return int(ev) if ev != "" else default
            except Exception:
                return default

if "_cfg_float" not in globals():
    def _cfg_float(name: str, default: float = 0.0) -> float:
        try:
            v = getattr(CFG, name)
            if v is None or str(v).strip() == "":
                raise AttributeError
            return float(str(v).strip())
        except Exception:
            ev = (os.getenv(name, "") or "").strip()
            try:
                return float(ev) if ev != "" else default
            except Exception:
                return default

if "_cfg_bool" not in globals():
    def _cfg_bool(name: str, default: bool = False) -> bool:
        try:
            v = getattr(CFG, name)
            if isinstance(v, bool):
                return v
            s = str(v).strip().lower()
            if s in ("1", "true", "yes", "on"):
                return True
            if s in ("0", "false", "no", "off"):
                return False
            return default
        except Exception:
            s = (os.getenv(name, "") or "").strip().lower()
            if s in ("1", "true", "yes", "on"):
                return True
            if s in ("0", "false", "no", "off"):
                return False
            return default

def _truthy_cfg(name: str, default: bool = False) -> bool:
    return _cfg_bool(name, default)

def _env_int(name: str, default: int) -> int:
    # (하위호환 alias) 기존 호출부 유지
    return _cfg_int(name, default)

def _env_float(name: str, default: float) -> float:
    # (하위호환 alias) 기존 호출부 유지
    return _cfg_float(name, default)

def ensure_config_fresh() -> None:
    """
    런타임에서 최신 .env/CFG 반영.
    - 기본: 프로세스당 최초 1회만 reload_config() 수행 (중복 로그 제거)
    - LOCAL_RAG_RELOAD_EACH_CALL=1 이면 매 호출 강제 허용(디버깅용)
    """
    global _RELOAD_ONCE_FLAG
    try:
        # 디버깅/테스트 강제 재로드 옵션
        force_each = (os.getenv("LOCAL_RAG_RELOAD_EACH_CALL", "").strip().lower()
                      in {"1", "true", "yes", "on"})
        if not force_each and _RELOAD_ONCE_FLAG:
            return
        with _RELOAD_LOCK:
            # 잠금 내에서 한 번 더 체크(경쟁 조건 방지)
            if not force_each and _RELOAD_ONCE_FLAG:
                return
            reload_config()
            _RELOAD_ONCE_FLAG = True
            # 최초 1회만 로그 남김 (중복 로그 억제)
            logger.debug("[LOCAL RAG] reload_config() applied (once).")
    except Exception:
        # 안전 무시
        pass

# ──────────────────────────────────────────────────────────────
# 화이트리스트(대용량 예외 허용) 유틸
#   - CFG.LOCAL_RAG_LARGE_WHITELIST: 콤마/세미콜론 구분 glob 패턴들
#       예) "*.국가보고서.pdf, *_final_report_*.pdf; KPI_*.xlsx"
#   - CFG.LOCAL_RAG_ALLOW_LARGE_PDF: bool, PDF만 전역 허용(선택)
#   - 패턴이 're:'로 시작하면 정규식으로 처리 (예: "re:.*(연차보고서|백서).*\\.pdf$")
# ──────────────────────────────────────────────────────────────
def _csv_patterns(val: str) -> list[str]:
    if not val:
        return []
    # 콤마/세미콜론/개행 구분을 모두 허용
    parts = re.split(r"[,\n;]+", val)
    return [p.strip() for p in parts if p and p.strip()]

def _is_large_whitelisted(path: str) -> bool:
    try:
        allow_pdf_all = _cfg_bool("LOCAL_RAG_ALLOW_LARGE_PDF", False)
        if allow_pdf_all and str(path).lower().endswith(".pdf"):
            return True
    except Exception:
        pass
    try:
        patterns = _csv_patterns(_cfg_str("LOCAL_RAG_LARGE_WHITELIST", ""))
        if not patterns:
            return False
        name = Path(path).name
        full = str(Path(path).resolve())
        for pat in patterns:
            if pat.startswith("re:"):
                rx = pat[3:]
                try:
                    if re.search(rx, name, flags=re.I) or re.search(rx, full, flags=re.I):
                        return True
                except Exception:
                    # 패턴 오류는 무시하고 계속
                    continue
            else:
                # glob(fnmatch) 는 파일명/전체경로 둘 다 시도
                if fnmatch(name, pat) or fnmatch(full, pat):
                    return True
        return False
    except Exception:
        return False

def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")

def _sha1_file(path: str) -> str:
    h = hashlib.sha1()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()

def _cache_dir() -> Path:
    cache_dir = _cfg_str("LOCAL_RAG_CACHE_DIR", "")
    if not cache_dir:
        cache_dir = str(Path(research_base_dir()) / ".cache")
    p = Path(cache_dir)
    p.mkdir(parents=True, exist_ok=True)
    return p

def _cache_paths(path: str) -> Tuple[str, str]:
    """대용량 변환 캐시 경로 결정: cache/{sha1}.{ext}.json"""
    cdir = _cache_dir()
    sha1 = _sha1_file(path)
    ext = Path(path).suffix.lower().lstrip(".")
    meta = (cdir / f"{sha1}.{ext}.meta.json").as_posix()
    data = (cdir / f"{sha1}.{ext}.data.json").as_posix()
    return meta, data

def _cache_load(path: str) -> Optional[List[dict]]:
    # CFG 우선 → ENV 폴백 (신규 일원화 헬퍼 사용)
    if _truthy_cfg("LOCAL_RAG_CACHE_IGNORE", False):
        return None
    meta, data = _cache_paths(path)
    try:
        if os.path.isfile(meta) and os.path.isfile(data):
            with open(meta, "r", encoding="utf-8") as fm:
                _ = json.load(fm)  # 향후 버전 검증 등
            with open(data, "r", encoding="utf-8") as fd:
                items = json.load(fd)
            logger.info("[LOCAL RAG] cache hit → %s", data)
            return items
    except Exception as e:
        logger.debug("cache load failed: %s", e)
    return None

def _cache_save(path: str, items_from_reader: List[dict]) -> None:
    try:
        meta, data = _cache_paths(path)
        with open(meta, "w", encoding="utf-8") as fm:
            json.dump({
                "path": str(Path(path).resolve()),
                "sha1": _sha1_file(path),
                "saved_at": _now_iso(),
                "count": len(items_from_reader),
            }, fm, ensure_ascii=False, indent=2)
        with open(data, "w", encoding="utf-8") as fd:
            json.dump(items_from_reader, fd, ensure_ascii=False, indent=2)
        logger.info("[LOCAL RAG] cache saved → %s", data)
    except Exception as e:
        logger.debug("cache save failed: %s", e)

def _file_version(path: str) -> str:
    """파일 변경을 반영하는 버전 식별자."""
    forced = _cfg_str("LOCAL_RAG_FORCE_VERSION", "")
    if forced:
        logger.debug("Using forced version for %s: %s", path, forced)
        return str(forced)

    mode = _cfg_str("LOCAL_RAG_VERSION_MODE", "mtime").lower()
    try:
        if mode == "sha1":
            h = hashlib.sha1()
            with open(path, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    h.update(chunk)
            ver = h.hexdigest()[:12]
        else:
            st = os.stat(path)
            ver = f"{int(st.st_mtime)}_{st.st_size}"
        logger.debug("Computed version for %s: %s (mode=%s)", path, ver, mode)
        return ver
    except Exception as e:
        logger.warning("Version compute failed for %s: %s", path, e)
        return "na"

def _truncate(s: str, max_chars_env: str = "LOCAL_RAG_MAX_TEXT_CHARS") -> str:
    # CFG 우선 → ENV
    cfg_val = getattr(CFG, max_chars_env, None)
    if cfg_val is None or str(cfg_val).strip() == "":
        max_chars = _cfg_int(max_chars_env, 200000)
    else:
        max_chars = int(str(cfg_val))
    if max_chars > 0 and len(s) > max_chars:
        return s[:max_chars]
    return s

# ──────────────────────────────────────────────────────────────
# 최소 길이 하한(게이트) — 너무 짧은 내용을 버리지 않되, 임계치 미만은 스킵
#   ENV/CFG 키(없으면 기본):
#     LOCAL_MIN_CHARS_GLOBAL=80, LOCAL_MIN_CHARS_PDF=60,
#     LOCAL_MIN_CHARS_PPTX=30,  LOCAL_MIN_CHARS_XLSX=1
# ──────────────────────────────────────────────────────────────
def _min_gate_for_ext(ext: str) -> int:
    e = (ext or "").lower()
    global_min = _cfg_int("LOCAL_MIN_CHARS_GLOBAL", 80)
    if e == ".pdf":
        return _cfg_int("LOCAL_MIN_CHARS_PDF", 60)
    if e == ".pptx":
        return _cfg_int("LOCAL_MIN_CHARS_PPTX", 30)
    if e == ".xlsx":
        return _cfg_int("LOCAL_MIN_CHARS_XLSX", 1)
    return global_min


# ──────────────────────────────────────────────────────────────
# 파일 우선순위 정렬키
#   findings.md > .pdf > .pptx > .xlsx > 기타
#   동일 우선순위에서는 (파일크기 내림차순, mtime 내림차순, 경로명 오름차순)
# ──────────────────────────────────────────────────────────────
_EXT_RANK: dict[str, int] = {
    # 파일타입 가중치: PDF(보고서/허가) > PPTX(리포트) > XLSX(정량)
    ".pdf": 1,
    ".pptx": 2,
    ".xlsx": 3,
}

def _is_findings_md(path: str) -> bool:
    name = Path(path).name.lower()
    # round-*-findings.md, *findings*.md, findings.md 모두 포착
    return name == "findings.md" or ("findings" in name and name.endswith(".md"))

def _path_priority(path: str) -> int:
    """
    경로/파일명에 기반한 우선순위(값이 작을수록 우선).
    예: refs/, 팩트북, 허가, 제품설명서 등을 상단으로.
    """
    p = str(path).lower()
    name = Path(path).name.lower()
    # 기본값(가장 뒤)
    pri = 10
    # 강한 신뢰 소스
    if "/refs/" in p or "\\refs\\" in p:
        pri = min(pri, 0)
    # 한글 키워드
    keywords = ["팩트북", "허가", "제품설명서"]
    if any(k in p for k in keywords) or any(k in name for k in keywords):
        pri = min(pri, 1)
    return pri


def _sort_key(path: str) -> tuple[int, int, int, int, str]:
    """
    반환 튜플이 작을수록 선순위가 되도록 구성.
    1) 경로/키워드 우선(_path_priority)
    2) findings.md 최상단
    3) 확장자 가중치 (.pdf < .pptx < .xlsx < 기타)
    4) 최근성 우선(mtime desc)
    5) 크기 내림차순
    """
    try:
        pri = _path_priority(path)
        if _is_findings_md(path):
            rank = -1  # findings를 최상단
        else:
            ext = Path(path).suffix.lower()
            rank = 1 + _EXT_RANK.get(ext, 99)  # 미지정 확장자는 가장 뒤
        st = os.stat(path)
        size = int(st.st_size)
        mtime = int(st.st_mtime)
    except Exception:
        # 접근 실패 시 최하위로 밀고 안전 키 반환
        pri, rank, size, mtime = 99, 999, 0, 0
    # (pri, rank, -mtime, -size, path)
    return (pri, rank, -mtime, -size, str(path).lower())


# ── Readers ──────────────────────────────────────────────────────────────────
def _read_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _read_md(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _read_csv(path: str) -> str:
    # 간단 CSV → 줄 단위 텍스트(필요시 unstructured/판다스로 확장)
    try:
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return Path(path).read_text(encoding="cp949", errors="ignore")

def _read_html(path: str) -> str:
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    if BeautifulSoup is not None:
        try:
            return BeautifulSoup(raw, "lxml").get_text(" ", strip=True)
        except Exception:
            return BeautifulSoup(raw, "html.parser").get_text(" ", strip=True)
    txt = re.sub(r"<[^>]+>", " ", raw)
    return re.sub(r"\s+", " ", txt).strip()

def _read_docx(path: str) -> str:
    if not docx:
        raise RuntimeError("python-docx 미설치")
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs)

def _read_pdf_pages(path: str, max_pages: Optional[int] = None) -> List[str]:
    """
    PDF 페이지 텍스트를 페이지별 리스트로 반환.
    max_pages 지정 없으면 ENV/CFG의 LOCAL_RAG_PDF_MAX_PAGES(기본 50) 사용.
    """
    if not PdfReader:
        raise RuntimeError("PyPDF2 미설치")
    reader = PdfReader(path)
    # 기본 페이지 제한 상향(보고서 중심): 50
    limit = int(max_pages) if (isinstance(max_pages, int) and max_pages > 0) else _env_int("LOCAL_RAG_PDF_MAX_PAGES", 50)
    pages: List[str] = []
    for i, p in enumerate(reader.pages, start=1):
        if limit and i > limit:
            break
        try:
            txt = p.extract_text() or ""
            pages.append(txt)
        except Exception:
            pages.append("")
    return pages

# ── PPTX helpers ─────────────────────────────────────────────────────────────
def _read_unstructured_elements(path: str) -> List[dict]:
    """
    PPTX, XLSX 등을 unstructured로 분해하여 요소(Element) 리스트로 반환.
    각 요소 dict 예시: {"sheet": "<슬라이드/시트/파트>", "row_index": <번호>, "content": "<텍스트>"}.
    """
    if partition is None:
        raise RuntimeError("unstructured 라이브러리 미설치")
    try:
        elements = partition(filename=path)
    except Exception as e:
        logger.warning("Unstructured partition failed for %s: %s", path, e)
        return []

    items: List[dict] = []
    for i, element in enumerate(elements):
        content = str(getattr(element, "text", "") or "").strip()
        if Table is not None and isinstance(element, Table):
            text_as_csv = getattr(element, "text_as_csv", "") or ""
            if text_as_csv:
                content = (content + "\n\n" + text_as_csv).strip()
        if not content:
            continue
        metadata = getattr(element, "metadata", None)
        part_name = getattr(metadata, "page_number", None) or getattr(metadata, "sheet_name", None) or "Part"
        items.append({"sheet": str(part_name), "row_index": i + 1, "content": content})
    return items

def _pick_sample_indices(n: int, ratio: float) -> List[int]:
    import math
    k = max(1, int(math.ceil(n * max(0.01, min(1.0, ratio)))))
    base = {0, max(0, n - 1)}
    if n <= len(base):
        return sorted(base)
    step = max(1, n // k)
    for i in range(0, n, step):
        base.add(i)
        if len(base) >= k:
            break
    return sorted(min(idx, n - 1) for idx in base)

def _pptx_extract_titles_bullets(path: str, sample_large: bool = False) -> List[dict]:
    """
    대용량 PPTX 전용: 제목+불릿만 추출, 표/차트 존재는 태그로 표기.
    sample_large=True면 슬라이드 일부 샘플링(기본 20%), 항상 1/마지막 포함.
    """
    if _PptxPresentation is None:
        raise RuntimeError("python-pptx 미설치")

    prs = _PptxPresentation(path)
    slides = list(prs.slides)
    total = len(slides)
    ratio = _env_float("LOCAL_RAG_SAMPLE_RATIO", 0.2) if sample_large else 1.0
    pick = _pick_sample_indices(total, ratio)

    items: List[dict] = []
    for sidx in pick:
        slide = slides[sidx]
        texts: List[str] = []

        # 제목 추출
        title = ""
        try:
            if getattr(slide.shapes, "title", None) is not None:
                title = str(getattr(slide.shapes.title, "text", "") or "").strip()
        except Exception:
            title = ""

        if title:
            texts.append(title)

        # 불릿 추출
        try:
            for shp in slide.shapes:
                has_tf = bool(getattr(shp, "has_text_frame", False))
                tf = getattr(shp, "text_frame", None)
                if not (has_tf and tf):
                    continue
                for p in getattr(tf, "paragraphs", []) or []:
                    raw = str(getattr(p, "text", "") or "").strip()
                    if not raw:
                        continue
                    if title and raw == title:
                        continue
                    level = int(getattr(p, "level", 0) or 0)
                    if level > 0 or len(raw) >= 2:
                        texts.append(f"- {raw}")
        except Exception:
            pass

        # 표/차트 존재 태그
        try:
            has_table = any(bool(getattr(shp, "has_table", False)) for shp in slide.shapes)
            has_chart = any(getattr(shp, "chart", None) is not None for shp in slide.shapes)
            if has_table:
                texts.append("[표 있음]")
            if has_chart:
                texts.append("[차트 있음]")
        except Exception:
            pass

        content = "\n".join(t for t in texts if t).strip()
        if not content:
            continue
        items.append({"sheet": str(sidx + 1), "row_index": sidx + 1, "content": content})

    logger.info("[LOCAL RAG] PPTX %s → slim extracted (slides=%d/%d)", Path(path).name, len(items), total)
    return items

def _pptx_extract_full(path: str, sample_large: bool = False) -> List[dict]:
    """
    PPTX 전체 추출: 제목/불릿 + 노트 + 표(→TSV)까지 텍스트화.
    sample_large=True면 일부 슬라이드 샘플링(기본 20%).
    """
    if _PptxPresentation is None:
        raise RuntimeError("python-pptx 미설치")
    prs = _PptxPresentation(path)
    slides = list(prs.slides)
    total = len(slides)
    ratio = _env_float("LOCAL_RAG_SAMPLE_RATIO", 0.2) if sample_large else 1.0
    pick = _pick_sample_indices(total, ratio)

    def _table_to_tsv(tbl) -> str:
        try:
            rows = []
            for r in tbl.rows:
                cells = []
                for c in r.cells:
                    txt = (getattr(c, "text", "") or "").strip()
                    cells.append(txt.replace("\t", " ").replace("\n", " "))
                rows.append("\t".join(cells))
            return "\n".join(rows)
        except Exception:
            return ""

    items: List[dict] = []
    for sidx in pick:
        slide = slides[sidx]
        buf: List[str] = []
        # 제목/불릿
        title = ""
        try:
            if getattr(slide.shapes, "title", None) is not None:
                title = str(getattr(slide.shapes.title, "text", "") or "").strip()
        except Exception:
            pass
        if title:
            buf.append(title)
        try:
            for shp in slide.shapes:
                if bool(getattr(shp, "has_text_frame", False)) and getattr(shp, "text_frame", None):
                    for p in getattr(shp.text_frame, "paragraphs", []) or []:
                        raw = (getattr(p, "text", "") or "").strip()
                        if raw and raw != title:
                            level = int(getattr(p, "level", 0) or 0)
                            buf.append(("- " if level > 0 else "") + raw)
                if bool(getattr(shp, "has_table", False)):
                    tsv = _table_to_tsv(getattr(shp, "table"))
                    if tsv.strip():
                        buf.append("[표]\n" + tsv)
        except Exception:
            pass
        # 노트
        try:
            ns = getattr(slide, "notes_slide", None)
            if ns and getattr(ns, "notes_text_frame", None):
                nt = (getattr(ns.notes_text_frame, "text", "") or "").strip()
                if nt:
                    buf.append("[노트]\n" + nt)
        except Exception:
            pass
        content = "\n".join(x for x in buf if x).strip()
        if content:
            items.append({"sheet": str(sidx + 1), "row_index": sidx + 1, "content": content})
    logger.info("[LOCAL RAG] PPTX %s → full extracted (slides=%d/%d)", Path(path).name, len(items), total)
    return items


def _read_pptx(path: str) -> List[dict]:
    """
    우선 python-pptx로 슬라이드 단위(제목+불릿) 추출 시도,
    실패/미설치 시 unstructured로 폴백.
    """
    if _truthy_cfg("SKIP_PPTX", False):
        logger.info("SKIP_PPTX=1 → _read_pptx skipped for %s", path)
        return []

    if _PptxPresentation is not None:
        try:
            return _pptx_extract_titles_bullets(path, sample_large=False)
        except Exception as e:
            logger.warning("python-pptx extract failed, fallback to unstructured: %s", e)

    if _PptxPresentation is not None:
        try:
            # 슬림/풀 선택: 기본은 풀 추출로 품질 우선, 대용량/샘플링은 슬림
            prefer_slim = _truthy_cfg("LOCAL_RAG_PPTX_SLIM", False)
            if prefer_slim:
                return _pptx_extract_titles_bullets(path, sample_large=False)
            return _pptx_extract_full(path, sample_large=False)
        except Exception as e:
            logger.warning("python-pptx extract failed, fallback to unstructured: %s", e)

    # fallback
    return _read_unstructured_elements(path)

def _read_xlsx_tsv(path: str, *, max_rows_per_sheet: Optional[int] = None) -> List[dict]:
    """
    XLSX을 시트별 TSV 텍스트로 직렬화. openpyxl 우선, 실패 시 unstructured 폴백.
    """
    rows_cap = max(0, int(max_rows_per_sheet or _cfg_int("LOCAL_RAG_XLSX_MAX_ROWS", 500)))
    if _OpenPyXL is None:
        logger.debug("openpyxl not installed; using unstructured for xlsx: %s", path)
        return _read_unstructured_elements(path)
    try:
        wb = _OpenPyXL.load_workbook(path, data_only=True, read_only=True)
    except Exception as e:
        logger.warning("openpyxl load failed, fallback to unstructured: %s", e)
        return _read_unstructured_elements(path)
    items: List[dict] = []
    for ws in wb.worksheets:
        lines: List[str] = []
        try:
            iter_rows = ws.iter_rows(values_only=True)
            cnt = 0
            for row in iter_rows:
                cnt += 1
                if rows_cap and cnt > rows_cap:
                    break
                cells = []
                for v in (row or ()):
                    s = "" if v is None else str(v)
                    s = s.replace("\t", " ").replace("\n", " ").strip()
                    cells.append(s)
                lines.append("\t".join(cells))
        except Exception:
            pass
        content = "\n".join([f"[시트] {ws.title}"] + lines).strip()
        if content:
            items.append({"sheet": ws.title, "row_index": 1, "content": content})
    logger.info("[LOCAL RAG] XLSX %s → tsv extracted (sheets=%d)", Path(path).name, len(items))
    return items


# ── Chunking (Markdown-like + Generic) ───────────────────────────────────────
_MD_FENCE_RE = re.compile(r"```.*?```", flags=re.DOTALL)
_MD_TABLE_LINE_RE = re.compile(r"^\s*\|.*\|\s*$")

def _summarize_table_lines(lines: List[str], max_rows: int = 3) -> str:
    head = [ln.strip() for ln in lines[:max_rows] if ln.strip()]
    if not head:
        return "[표 요약]"
    sample = " / ".join(head)
    return f"[표 요약: {sample[:120]}...]"

def _base_chunk_params_from_env() -> tuple[int, int]:
    """
    기본 청크 크기/오버랩: 로컬 전용 상향 기본값.
    - LOCAL_CHUNK_SIZE(기본 1000)
    - LOCAL_CHUNK_OVERLAP(기본 120)
    """
    size = _cfg_int("LOCAL_CHUNK_SIZE", 1000)
    overlap = _cfg_int("LOCAL_CHUNK_OVERLAP", 120)
    return max(200, size), max(0, min(overlap, size // 2))

def _type_chunk_params(ext: str) -> tuple[int, int, int, str]:
    """
    파일 확장자별 (min_chars, max_chars, overlap, mode) 반환.
    mode: 'paragraph' | 'lines'
    """
    size, overlap = _base_chunk_params_from_env()
    e = (ext or "").lower()
    # 기본 min/max는 size를 중심으로 ±15%
    def around(sz: int, pct: float = 0.15) -> tuple[int, int]:
        mn = int(sz * (1.0 - pct))
        mx = int(sz * (1.0 + pct))
        return max(300, mn), max(mn + 200, mx)

    # 타입별 튜닝
    if e == ".pdf":
        mn, mx = around(max(800, size))
        return mn, max(mx, 1200), max(overlap, 100), "paragraph"
    if e == ".pptx":
        mn = _cfg_int("MIN_CHUNK_PPTX", 80)
        mx = 1200
        return mn, mx, 0, "paragraph"  # lines → paragraph, overlap=0
    if e in (".xlsx", ".csv"):
        mn, mx = 100, 1100  # 표/셀 데이터 특성상 짧을 수 있음
        return mn, mx, max(80, min(overlap, 120)), "paragraph"
    if e in (".md", ".markdown"):
        mn, mx = 100, 1300
        return mn, mx, max(overlap, 120), "paragraph"
    mn, mx = 100, 1300     # txt/html/docx 등
    return mn, mx, max(overlap, 120), "paragraph"

def _split_text_generic(
    text: str,
    *,
    min_chars: int,
    max_chars: int,
    overlap: int = 0,
    mode: str = "paragraph",
    markdown_features: bool = False,
) -> List[str]:
    """
    범용 스플리터:
      - mode='paragraph': 빈줄 기준 문단 병합/분할
      - mode='lines'    : 줄 단위(슬라이드·표 위주)에 친화적
      - markdown_features=True이면 MD 코드펜스/테이블 요약 적용
      - 오버랩(overlap) 지원: 청크 경계에서 뒤쪽 꼬리를 다음 청크 앞에 덧붙임
    """
    if not text:
        return []

    if markdown_features:
        # 1) 코드블록 요약 치환
        text = _MD_FENCE_RE.sub("[코드 요약]", text)
        # 2) 테이블 블록 요약
        lines = text.splitlines()
        out_lines: List[str] = []
        buf_table: List[str] = []
        def _flush_table():
            nonlocal buf_table, out_lines
            if buf_table:
                out_lines.append(_summarize_table_lines(buf_table))
                buf_table = []
        for ln in lines:
            if _MD_TABLE_LINE_RE.match(ln):
                buf_table.append(ln)
            else:
                _flush_table()
                out_lines.append(ln)
        _flush_table()
        text = "\n".join(out_lines)

    # 3) 기본 라인/문단 분해
    if mode == "lines":
        units = [ln.strip() for ln in text.splitlines() if ln.strip()]
    else:
        units = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    # 4) 병합(짧은 단위 → 다음과 병합)
    merged, buf = [], ""
    for u in units:
        if len(buf) + len(u) + 1 < min_chars:
            buf = (buf + " " + u).strip() if buf else u
        else:
            if buf:
                merged.append(buf)
            buf = u
    if buf:
        merged.append(buf)

    # 5) 최대 길이 초과 분절(단어 경계 우선)
    chunks: List[str] = []
    for m in merged:
        cur = m.strip()
        while len(cur) > max_chars:
            cut = cur[:max_chars]
            cut_idx = cut.rfind(" ")
            if cut_idx < max_chars * 0.5:
                cut_idx = max_chars
            head = cur[:cut_idx].strip()
            tail = cur[cut_idx:].strip()
            if head:
                chunks.append(head)
            cur = tail
        if cur:
            chunks.append(cur)

    if overlap <= 0 or len(chunks) <= 1:
        return chunks

    # 6) 오버랩 적용: 이전 청크 끝 꼬리를 다음 청크 앞에 덧붙임
    overlapped: List[str] = []
    prev_tail = ""
    for i, ch in enumerate(chunks):
        if i == 0:
            overlapped.append(ch)
            prev_tail = ch[-overlap:]
        else:
            prefix = (prev_tail + " ") if prev_tail else ""
            merged_ch = (prefix + ch).strip()
            overlapped.append(merged_ch)
            prev_tail = ch[-overlap:]
    return overlapped


def _ensure_min_chunk(content: str, ext: str, chunks: List[str]) -> List[str]:
    """
    청크가 비었지만 내용이 최소 하한을 넘으면 단일 청크로 수용.
    """
    if chunks:
        return chunks
    gate = _min_gate_for_ext(ext)
    txt = (content or "").strip()
    if len(txt) >= max(0, gate):
        return [txt]
    return []

def _file_uri(path: str) -> str:
    return Path(path).resolve().as_uri()  # file:/// 형태 보장

# ──────────────────────────────────────────────────────────────
# URL 안정화: stored_urls/seen-hash, dedupe를 위해 url은 base로 고정
#  - fragment/query를 url에 섞지 않고 locator로 분리한다.
# ──────────────────────────────────────────────────────────────
def _make_locator(*, part_label: str, index_num: Any, chunk: int, kind: str) -> str:
    """
    UI/디버깅용 위치 힌트.
    url에는 fragment를 넣지 않고, 별도 locator로 보관한다.
    """
    try:
        return f"{kind}:{part_label}:{index_num}:chunk:{int(chunk)}"
    except Exception:
        return f"{kind}:{part_label}:{index_num}:chunk:{chunk}"


# ──────────────────────────────────────────────────────────────
# part 메타 구성 유틸: kind와 인덱스를 안정 문자열로 결합
# ──────────────────────────────────────────────────────────────
def _compose_part(kind: str, idx: Any) -> str:
    try:
        k = (kind or "").strip().lower()
    except Exception:
        k = "part"
    try:
        v = str(idx).strip()
    except Exception:
        v = ""
    return f"{k}:{v}" if v else k


# ──────────────────────────────────────────────────────────────────────────────
# 변환: 로컬 파일 → web.json 아이템 배열
# ──────────────────────────────────────────────────────────────────────────────
def _to_webjson_items(path: str, *, max_pages_per_file: Optional[int] = None) -> List[dict]:
    ext = Path(path).suffix.lower()
    title = Path(path).name
    url_click = _file_uri(path)  # 사람이 눌러 열어볼 주소
    ver = _file_version(path)
    fetched_at = _now_iso()
    # source_version 안정화용: 파일 mtime 포함 (web_rag SSoT와 동일하게 사용)
    try:
        _st = os.stat(path)
        _mtime = int(_st.st_mtime)
    except Exception:
        _mtime = 0

    _ct_map = {
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".html": "text/html",
        ".htm": "text/html",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv": "text/csv",
    }
    content_type = _ct_map.get(ext, "application/octet-stream")

    # 용량 가드(+ 대용량 PPTX 슬림/샘플링)
    max_mb = _cfg_float("LOCAL_RAG_MAX_FILE_MB", 50.0)
    try:
        _filesize_bytes = os.path.getsize(path)
        size_mb = _filesize_bytes / (1024 * 1024)
    except Exception:
        _filesize_bytes = 0
        size_mb = None
    is_large = bool(size_mb and max_mb > 0 and size_mb > max_mb)
    enable_sample = _truthy_cfg("LOCAL_RAG_SAMPLE_LARGE", False)
    # 대용량 파일 예외 허용(화이트리스트)
    large_ok = _is_large_whitelisted(path)

    text: Optional[str] = None
    items_from_reader: Optional[List[dict]] = None

    try:
        if ext == ".pdf":
            if is_large and not large_ok:
                logger.info("[LOCAL RAG] skip large PDF (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            if is_large and large_ok:
                logger.info("[LOCAL RAG] WHITELISTED large PDF allowed: %s (%.1f MB)", path, size_mb or -1)
            pages = _read_pdf_pages(path, max_pages=max_pages_per_file)
            items_from_reader = []
            for i, txt in enumerate(pages, start=1):
                if (txt or "").strip():
                    items_from_reader.append({"page_num": i, "content": txt})

        elif ext in (".txt", ".md", ".markdown", ".html", ".htm", ".docx", ".csv"):
            if is_large and not large_ok:
                logger.info("[LOCAL RAG] skip large text-like (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            if is_large and large_ok:
                logger.info("[LOCAL RAG] WHITELISTED large text-like allowed: %s (%.1f MB)", path, size_mb or -1)
            if ext == ".txt":
                text = _read_txt(path)
            elif ext in (".md", ".markdown"):
                text = _read_md(path)
            elif ext in (".html", ".htm"):
                text = _read_html(path)
            elif ext == ".docx":
                text = _read_docx(path)
            elif ext == ".csv":
                text = _read_csv(path)

        elif ext == ".pptx":
            if is_large and large_ok:
                logger.info("[LOCAL RAG] WHITELISTED large PPTX allowed: %s (%.1f MB)", path, size_mb or -1)
                # 화이트리스트라도 기본은 슬림 추출로 부하 경감
                prefer_slim = True
                items_from_reader = _pptx_extract_titles_bullets(path, sample_large=True)
                _cache_save(path, items_from_reader or [])
            elif is_large and enable_sample:
                cached = _cache_load(path)
                if cached is not None:
                    items_from_reader = cached
                    logger.info("[LOCAL RAG] sampled/slim extracted (cache) → %s", Path(path).name)
                else:
                    items_from_reader = _pptx_extract_titles_bullets(path, sample_large=True)
                    _cache_save(path, items_from_reader or [])
                    logger.info("[LOCAL RAG] sampled/slim extracted (fresh) → %s", Path(path).name)
            elif is_large and not enable_sample:
                logger.info("[LOCAL RAG] skip large file (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            else:
                prefer_slim = _truthy_cfg("LOCAL_RAG_PPTX_SLIM", False)
                if _PptxPresentation is not None:
                    items_from_reader = (_pptx_extract_titles_bullets if prefer_slim else _pptx_extract_full)(path, sample_large=False)
                else:
                    items_from_reader = _read_unstructured_elements(path)

        elif ext == ".xlsx":
            if is_large and not large_ok:
                logger.info("[LOCAL RAG] skip large XLSX (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            if is_large and large_ok:
                logger.info("[LOCAL RAG] WHITELISTED large XLSX allowed: %s (%.1f MB)", path, size_mb or -1)
            items_from_reader = _read_xlsx_tsv(path)

        else:
            logger.debug("[LOCAL RAG] unsupported extension skipped: %s", path)
            return []

    except RuntimeError as e:
        logger.warning("[LOCAL RAG] Reader dependency failed (RuntimeError): %s -> %s", path, e)
        return []
    except Exception as e:
        logger.warning("[LOCAL RAG] Extraction error: %s -> %s", path, e)

    # 2) 결과 통합 및 메타데이터 생성 (공통)
    final_items: List[dict] = []
    _pri = _path_priority(path)  # 아이템 레벨 정렬에도 활용

    # 리스트 형태 결과(PDF, PPTX, XLSX)
    if items_from_reader is not None and isinstance(items_from_reader, list):
        # 타입별 청크 파라미터
        minc, maxc, overlap, mode = _type_chunk_params(ext)
        for it in items_from_reader:
            # kind/label/index 추출 (pop 이전에 미리 읽음)
            has_page   = "page_num" in it
            has_sheet  = "sheet" in it
            has_row    = "row_index" in it
            # 파일 확장자 기반 기본 kind
            if ext == ".pdf":
                kind = "page"
            elif ext == ".pptx":
                kind = "slide"
            elif ext == ".xlsx":
                kind = "sheet"
            else:
                kind = "part"
            # label/index 결정
            label_raw = it.get("sheet") if has_sheet else (it.get("page_num") if has_page else "Part")
            idx_raw   = it.get("row_index") if has_row else (it.get("page_num") if has_page else 1)
            # 기존 로직 유지(pop)하되 위에서 읽은 값 사용
            part_label = str(it.pop("sheet", it.pop("page_num", label_raw)))
            index_num  = it.pop("row_index", it.pop("page_num", idx_raw))
            content = (it.pop("content", "") or "").strip()
            if not content:
                continue
            chunks = _split_text_generic(
                content,
                min_chars=minc,
                max_chars=maxc,
                overlap=overlap,
                mode=("lines" if mode == "lines" else "paragraph"),
                markdown_features=(ext in (".md", ".markdown")),
            )
            if not chunks:
                continue
            chunks = _ensure_min_chunk(content, ext, chunks)
            if not chunks:
                continue
            for j, ch in enumerate(chunks, start=1):
                part_value = _compose_part(kind, f"{part_label}:{index_num}")
                # SSoT: source is canonical file:// URI + fragment (part/index/chunk)
                # ingest_vector will generate ids from (source, source_version, content) using utils.make_doc_id
                _src = _build_local_source(Path(path), part=part_label, index=index_num, chunk=j)
                _ver = str(_mtime)
                final_items.append({
                    "title": f"{title} ({part_label}, Index: {index_num}, Chunk {j})",
                    "url": url_click,
                    "source": _src,
                    "source_version": _ver,
                    "part": part_value,
                    "locator": _make_locator(part_label=str(part_label), index_num=index_num, chunk=j, kind=kind),
                    "content": ch,
                    "content_type": content_type,
                    "bytes": _filesize_bytes,
                    "fetched_at": fetched_at,
                    "mtime": _mtime,
                    "pri": _pri,
                })

    # 단일 텍스트 결과(TXT, MD, HTML, DOCX, CSV)
    elif text is not None:
        content = _truncate(text or "", "LOCAL_RAG_MAX_TEXT_CHARS")
        if content.strip():
            minc, maxc, overlap, mode = _type_chunk_params(ext)
            chunks = _split_text_generic(
                content,
                min_chars=minc,
                max_chars=maxc,
                overlap=overlap,
                mode=("lines" if mode == "lines" else "paragraph"),
                markdown_features=(ext in (".md", ".markdown")),
            )
            chunks = _ensure_min_chunk(content, ext, chunks)
            for j, ch in enumerate(chunks, start=1):
                part_value = _compose_part("chunk", j)
                _src = _build_local_source(Path(path), part="chunk", index=None, chunk=j)
                _ver = str(_mtime)
                final_items.append({
                    "title": f"{title} (Chunk {j})",
                    "url": url_click,
                    "source": _src,
                    "source_version": _ver,
                    "part": part_value,
                    "locator": _make_locator(part_label="chunk", index_num=j, chunk=j, kind="chunk"),
                    "content": ch,
                    "content_type": content_type,
                    "bytes": _filesize_bytes,
                    "fetched_at": fetched_at,
                    "mtime": _mtime,
                    "pri": _pri,
                })

    if not final_items:
        logger.debug("[LOCAL RAG] File yielded no extractable content: %s", path)
        return []

    return final_items

# ──────────────────────────────────────────────────────────────────────────────
# 엔트리: globs → web.json 생성
# ──────────────────────────────────────────────────────────────────────────────
def build_webjson_from_local(
    globs: List[str],
    out_dir: str,
    *,
    max_docs: Optional[int] = None,
    max_pages_per_file: Optional[int] = None,
) -> str:
    """
    globs에 매칭되는 로컬 파일들을 읽어 web.json으로 변환.
    - max_docs: 전체 아이템(청크) 상한. None/0이면 무제한. ENV: LOCAL_RAG_MAX_DOCS
    - max_pages_per_file: 파일별 페이지 상한(PDF 등). None/0이면 무제한. ENV: LOCAL_RAG_MAX_PAGES_PER_FILE
    """
    # 런타임 설정 갱신
    ensure_config_fresh()
    # CFG 우선 → ENV 폴백
    if max_docs is None:
        max_docs = _cfg_int("LOCAL_RAG_MAX_DOCS", 0)
    if max_pages_per_file is None:
        # 기본 페이지 캡 상향(보고서 중심): 50
        max_pages_per_file = _cfg_int("LOCAL_RAG_MAX_PAGES_PER_FILE", 50)

    logger.info("[LOCAL RAG] CWD: %s", os.getcwd())
    logger.info("[LOCAL RAG] Received globs: %s", globs)
    if max_docs:
        logger.info("[LOCAL RAG] cap: max_docs=%d", max_docs)
    if max_pages_per_file:
        logger.info("[LOCAL RAG] cap: max_pages_per_file=%d", max_pages_per_file)

    Path(out_dir).mkdir(parents=True, exist_ok=True)
    files: List[str] = []

    for g in globs or []:
        g = os.path.expandvars(os.path.expanduser(g))
        logger.debug("[LOCAL RAG] Expanded glob pattern: %s", g)
        matched = glob.glob(g, recursive=True)
        files.extend(matched)
        logger.info("[LOCAL RAG] Pattern %s matched %d files.", g, len(matched))

    files = sorted({f for f in files if os.path.isfile(f)})
    logger.info("[LOCAL RAG] Total unique files found: %d", len(files))

    # [P1-1] 우선순위 정렬: findings.md > PDF > PPTX > XLSX > 기타
    files.sort(key=_sort_key)

    # [P1-1] cap 적용은 정렬 이후: LOCAL_RAG_MAX_FILES(우선) → LOCAL_MAX_FILES(폴백) → 기본 1500
    _cap = _cfg_int("LOCAL_RAG_MAX_FILES", _cfg_int("LOCAL_MAX_FILES", 1500))
    if _cap > 0 and len(files) > _cap:
        # cap 적용 전/후 타입 분포를 로깅하여 신뢰도 향상
        def _ext_tag(p: str) -> str:
            if _is_findings_md(p):
                return "findings.md"
            return Path(p).suffix.lower() or "(noext)"
        before = Counter(_ext_tag(p) for p in files)
        files = files[:_cap]
        after = Counter(_ext_tag(p) for p in files)
        logger.info("[LOCAL RAG] applying file cap (after priority sort): %d → %d", sum(before.values()), _cap)
        try:
            _b = ", ".join(f"{k}:{v}" for k, v in before.most_common())
            _a = ", ".join(f"{k}:{v}" for k, v in after.most_common())
            logger.info("[LOCAL RAG] file mix before cap: %s", _b)
            logger.info("[LOCAL RAG] file mix  after cap: %s", _a)
        except Exception:
            pass


    items: List[dict] = []
    processed_items = 0
    # cap은 전역 정렬 후 적용하므로 루프 중간에 절단하지 않습니다.
    total_files = len(files)

    for fi, f in enumerate(files, start=1):
        try:
            # 파일별 파싱 캐시: mtime+size가 같으면 재파싱 스킵
            cached = _cache_load(f)
            if cached is not None:
                file_items = cached
            else:
                file_items = _to_webjson_items(f, max_pages_per_file=max_pages_per_file or None)
                if file_items:
                    _cache_save(f, file_items)
        except Exception as e:
            logger.warning("[LOCAL RAG] local ingest 실패: %s -> %s", f, e)
            file_items = []

        for it in file_items:
            # 빈 content 제거
            if not (it.get("content") or "").strip():
                continue
            items.append(it)
            processed_items += 1

            # 500개 단위 진행률 로그
            if processed_items % 500 == 0:
                logger.info("[LOCAL RAG] processed %d items (files %d/%d)", processed_items, fi, total_files)

        # 루프 내 조기 종료 제거: 전 파일 수집 후 전역 정렬/절단
        pass

    # 중복 source 제거(최초 등장 우선)
    dedup_seen, dedup_items = set(), []
    for it in items:
        sid = it.get("source")
        if not sid or sid in dedup_seen:
            continue
        dedup_seen.add(sid)
        dedup_items.append(it)
    items = dedup_items

    # ──────────────────────────────────────────────────────────
    # [유지] 전역 아이템 우선순위 정렬 후 max_docs 적용
    #       (파일 선별은 이미 위에서 완료했으며, 아이템 레벨에서도 타입/크기 기준을 유지)
    # ──────────────────────────────────────────────────────────
    # 아이템 정렬 가중치: PDF > PPTX > XLSX  (※ 아래에서 -_rank 로 정렬하므로 값이 클수록 우선)
    _rank: dict[str, int] = {
        "application/pdf": 3,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": 2,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": 1,
    }
    try:
        items.sort(
            key=lambda x: (
                -_rank.get((x.get("content_type") or ""), 0),        # 타입 가중치 (desc)
                int(x.get("pri", 10) or 10),                         # 경로 우선(asc)
                -int(x.get("mtime", 0) or 0),                        # 최근성(desc)
                -int(x.get("bytes", 0) or 0),                        # 크기(desc)
                str(x.get("source") or ""),
            )
        )
    except Exception as e:
        logger.debug("[LOCAL RAG] sort by priority skipped: %s", e)

    if max_docs and max_docs > 0 and len(items) > max_docs:
        logger.info("[LOCAL RAG] apply max_docs after sort: %d → %d", len(items), max_docs)
        items = items[:max_docs]


    # 디버그 샘플
    uniq_sources = {it.get("source") for it in items}
    logger.info("[LOCAL RAG] files=%d items=%d unique_sources=%d", len(files), len(items), len(uniq_sources))
    if items:
        def _pretty_src(src: str) -> str:
            s = unquote(src or "")
            # NOTE:
            # local sources are now canonical file:// URIs with optional fragment.
            # legacy "__v_" suffix stripping is no longer applicable.
            if "__p_" in s:
                s = s.split("__p_")[0]
            return s

        sample = items[:3]
        sample_titles = [it.get("title", "") for it in sample]
        sample_sources = [_pretty_src(it.get("source", "")) for it in sample]
        sample_urls = [unquote(it.get("url", "")) for it in sample]

        logger.debug("[LOCAL RAG] sample titles : %s", sample_titles)
        logger.debug("[LOCAL RAG] sample sources: %s", sample_sources)
        logger.debug("[LOCAL RAG] sample urls   : %s", sample_urls)

    # 헬스체크: 아이템 0개면 설정에 따라 중지 또는 경고
    if not items:
        if _cfg_bool("LOCAL_RAG_ALLOW_EMPTY", False):
            logger.warning("[LOCAL RAG] 0 items produced — continuing (LOCAL_RAG_ALLOW_EMPTY=1)")
            from datetime import datetime as _dt
            _ts = _dt.now().strftime("%Y_%m_%d_%H%M%S")
            out_path = os.path.join(out_dir, f"local_{_ts}.json")
            with open(out_path, "w", encoding="utf-8") as fp:
                json.dump([], fp)
            return out_path
        raise RuntimeError("LOCAL RAG produced 0 items after extraction/filters; aborting to prevent empty index.")

    ts = datetime.now().strftime("%Y_%m_%d_%H%M%S")
    out_path = os.path.join(out_dir, f"local_{ts}.json")
    with open(out_path, "w", encoding="utf-8") as fp:
        json.dump(items, fp, ensure_ascii=False, indent=2)
    logger.info("[LOCAL RAG] web.json saved → %s", out_path)
    return out_path

# ──────────────────────────────────────────────────────────────────────────────
# 파이프: web.json → Chroma 적재 + 미리보기
# ──────────────────────────────────────────────────────────────────────────────
def ingest_local_files(
    globs: List[str],
    namespace: str,
    persist_directory: str | None,
    topic_slug: str,
    root_dir: str,
    add_web_pages_json_to_chroma=None,
    web_page_json_to_documents=None,
) -> Tuple[List[str], List[Document], int]:
    """
    반환: (생성 JSON 경로들, 미리보기 Documents, 인덱싱된 청크 수 합)
    """
    # 런타임 설정 갱신
    ensure_config_fresh()
    # 연구 요약(findings) 포함 옵션
    include_findings = _truthy_cfg("INCLUDE_FINDINGS_IN_VECTOR", False)
    if include_findings:
        slug = (topic_slug or os.getenv("TOPIC_SLUG") or "default").strip()
        findings_dir = research_topic_dir(slug)
        findings_pattern = str(findings_dir / "round-*-findings.md")
        globs = list(globs or [])
        globs.append(findings_pattern)
        logger.info("[LOCAL RAG] findings included → %s", findings_pattern)

        # [ADD] 추가적인 findings 패턴(선택)
        extra_findings = [
            str(findings_dir / "*findings*.md"),
            str(findings_dir / "findings.md"),
        ]
        for pat in extra_findings:
            if pat not in globs:
                globs.append(pat)
        logger.debug("[LOCAL RAG] extra findings patterns included: %s", extra_findings)

    if not globs:
        logger.info("[LOCAL RAG] no globs provided → skip ingest")
        return ([], [], 0)

    # 표준 리소스 디렉터리 정책 사용
    res_dir = str(research_resources_dir(topic_slug or _cfg_str("TOPIC_SLUG", "default")))
    json_path = build_webjson_from_local(globs, res_dir)

    # ──────────────────────────────────────────────────────────
    # [ADD] web.json 기반 로컬 청크 통계 로깅
    # ──────────────────────────────────────────────────────────
    def _load_items_safe(p: str) -> list[dict]:
        try:
            with open(p, "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, list):
                return [x for x in data if isinstance(x, dict)]
        except Exception as e:
            logger.debug("[LOCAL RAG][stats] load failed for %s: %s", p, e)
        return []

    def _percentile(sorted_vals: List[int], q: float) -> int:
        if not sorted_vals:
            return 0
        q = max(0.0, min(1.0, q))
        idx = int(round((len(sorted_vals) - 1) * q))
        return int(sorted_vals[idx])

    try:
        _items_for_stats = _load_items_safe(json_path)
        _lens: List[int] = []
        _by_ct: Dict[str, int] = {}
        for it in _items_for_stats:
            c = (it.get("content") or "").strip()
            if not c:
                continue
            _lens.append(len(c))
            ct = (it.get("content_type") or "").strip() or "unknown"
            _by_ct[ct] = _by_ct.get(ct, 0) + 1
        _n = len(_lens)
        if _n:
            _lens.sort()
            _avg = sum(_lens) / _n
            _p50 = _percentile(_lens, 0.50)
            _p90 = _percentile(_lens, 0.90)
            _minv, _maxv = _lens[0], _lens[-1]
            logger.info("[LOCAL RAG][stats] items=%d avg_chars=%.1f p50=%d p90=%d min=%d max=%d",
                        _n, _avg, _p50, _p90, _minv, _maxv)
            if _by_ct:
                _top_ct = ", ".join(f"{k}:{v}" for k, v in sorted(_by_ct.items(), key=lambda kv: (-kv[1], kv[0]))[:6])
                logger.info("[LOCAL RAG][stats] by_content_type: %s", _top_ct)
            # (선택) 메트릭 이벤트 훅
            try:
                from tools.metrics import event as _metrics_event  # 안전 임포트
                _metrics_event("local_ingest_stats",
                               items=_n,
                               avg_chars=round(_avg, 1),
                               p50=_p50, p90=_p90,
                               min_chars=_minv, max_chars=_maxv,
                               by_ct=_by_ct)
            except Exception:
                pass
        else:
            logger.info("[LOCAL RAG][stats] items=0 (no content)")
    except Exception as e:
        logger.debug("[LOCAL RAG][stats] compute failed: %s", e)

    # ──────────────────────────────────────────────────────────
    # [ADD] 라운드 분할 인제스트 (MAX_LOCAL_DOCS_PER_ROUND)
    # ──────────────────────────────────────────────────────────
    chunk_total = 0
    out_json_paths: List[str] = [json_path]
    try:
        max_per_round = _cfg_int("MAX_LOCAL_DOCS_PER_ROUND", 0)
        if max_per_round and max_per_round > 0:
            with open(json_path, "r", encoding="utf-8") as f:
                _items_all = json.load(f)
            if isinstance(_items_all, list) and len(_items_all) > max_per_round:
                # 안정적 우선순위를 유지하기 위해 현재 순서 유지(= 파일 우선순위 반영된 상태)
                rounds: List[List[dict]] = [
                    _items_all[i:i + max_per_round] for i in range(0, len(_items_all), max_per_round)
                ]
                out_json_paths = []
                for ri, arr in enumerate(rounds, start=1):
                    rpath = os.path.join(os.path.dirname(json_path),
                                         f"{Path(json_path).stem}_r{ri}{Path(json_path).suffix}")
                    with open(rpath, "w", encoding="utf-8") as rf:
                        json.dump(arr, rf, ensure_ascii=False, indent=2)
                    out_json_paths.append(rpath)
                logger.info("[LOCAL RAG] split web.json into %d rounds (max_per_round=%d)", len(out_json_paths), max_per_round)
    except Exception as e:
        logger.debug("[LOCAL RAG] round-split skipped: %s", e)

    # ──────────────────────────────────────────────────────────
    # [NEW] 네임스페이스 정규화 + persist_directory 해석(ingest/retrieve 통일)
    # ──────────────────────────────────────────────────────────
    def _compute_effective_ns(ns_in: Optional[str], slug: str) -> str:
        """
        규칙:
          1) 인자 ns_in 우선
          2) CFG.CHROMA_NAMESPACE_LOCAL (없으면)
          3) f"{slug}-local" (없으면 "default-local")
        → 최종적으로 utils._sanitize_ns 규칙으로 정규화
        """
        base = (ns_in or
                getattr(CFG, "CHROMA_NAMESPACE_LOCAL", None) or
                (f"{(slug or 'default').strip() or 'default'}-local"))
        # Optional[Callable]이므로 None 체크 + callable 체크 모두 수행
        if _wr_sanitize_ns is not None:
            try:
                fn = _wr_sanitize_ns
                if callable(fn):
                    return fn(base)
            except Exception:
                pass
        # 폴백(간단 정규화): 영숫자/.-_ 외 '_'로 치환
        import re as _re
        s = (base or "").strip()
        s = _re.sub(r"[\\/]+", "_", s)
        s = _re.sub(r"[^A-Za-z0-9._\-]+", "_", s)
        s = s.strip("._-").lower() or "ns_default"
        return s

    def _compute_persist_dir(ns_eff: str, pd_in: Optional[str]) -> Optional[str]:
        """
        utils._resolve_persist_dir와 동일 규칙으로 디렉터리 결정.
        입력 persist_directory가 주어졌다면 그대로 사용(공백 제외),
        없으면 CHROMA_DIR 또는 DATA_DIR/chroma_store/<ns>.
        """
        if pd_in is None:
            fn2 = _wr_resolve_persist_dir
            if fn2 is not None and callable(fn2):
                try:
                    return fn2(ns_eff, None)
                except Exception:
                    pass
            # 폴백: 간단 구현 (utils와 동일 로직 요약)
            chroma_dir = str(getattr(CFG, "CHROMA_DIR", "") or "").strip()
            if chroma_dir:
                p = Path(chroma_dir)
                if p.name == ns_eff:
                    return str(p)
                if p.parent.name == "chroma_store":
                    return str(p.parent / ns_eff)
                return str(p / ns_eff)
            # DATA_DIR/chroma_store/<ns>
            _cand = os.path.join(str(research_base_dir()), "..")  # 프로젝트 루트 근처 폴백 시도
            try:
                from tools.web_rag.utils import DATA_DIR as _wr_DATA_DIR  # 최선
                return str(Path(_wr_DATA_DIR) / "chroma_store" / ns_eff)
            except Exception:
                return str(Path(_cand).resolve() / "data" / "chroma_store" / ns_eff)
        s = (pd_in or "").strip()
        return s or None

    ns_eff = _compute_effective_ns(namespace, topic_slug or _cfg_str("TOPIC_SLUG", "default"))
    pd_eff = _compute_persist_dir(ns_eff, persist_directory)

    if add_web_pages_json_to_chroma is not None:
        for rjson in out_json_paths:
            try:
                _orig, chunk_count = add_web_pages_json_to_chroma(
                    rjson, namespace=ns_eff, persist_directory=pd_eff
                )
                chunk_total += int(chunk_count or 0)
                logger.info("[LOCAL RAG] added to chroma: %s → chunks=%s (ns=%s, dir=%s)",
                            rjson, chunk_count, ns_eff, pd_eff)
            except Exception as e:
                logger.warning("[LOCAL RAG] add_web_pages_json_to_chroma(local) 실패(%s): %s", rjson, e)

    docs_preview: List[Document] = []
    if web_page_json_to_documents is not None:
        try:
            docs_preview = web_page_json_to_documents(json_path)[:8]
            logger.debug("[LOCAL RAG] preview docs: %d", len(docs_preview))
        except Exception as e:
            logger.warning("[LOCAL RAG] preview build(local) 실패: %s", e)

    return (out_json_paths, docs_preview, chunk_total)


# ──────────────────────────────────────────────────────────────────────────────
# [ADD] findings 빠른 인제스트: routers.after_synth에서 호출
# ──────────────────────────────────────────────────────────────────────────────
def add_local_findings_to_chroma(
    topic_slug: str,
    *,
    patterns: Optional[list[str]] = None,
    namespace: Optional[str] = None,
    persist_directory: Optional[str] = None,
) -> int:
    """
    방금 생성된 findings(md)를 신속히 Chroma에 반영.
    반환: 인덱싱된 청크 수(합계)
    """
    ensure_config_fresh()
    slug = (topic_slug or _cfg_str("TOPIC_SLUG", "default")).strip() or "default"

    # 기본 패턴: 프로젝트 표준 research/<topic_slug>/ 경로 가정
    if not patterns:
        base = research_topic_dir(slug)
        patterns = [
            str(base / "round-*-findings.md"),
            str(base / "*findings*.md"),
            str(base / "findings.md"),
        ]

    # 지연 임포트로 순환 의존 회피
    _add_to_chroma = None
    _to_docs = None
    try:
        from tools.web_rag import add_web_pages_json_to_chroma as _aw
        from tools.web_rag import web_page_json_to_documents as _wpj2d
        _add_to_chroma = _aw
        _to_docs = _wpj2d
    except Exception:
        pass

    # ingest_local_files 재사용
    #  - ingest_local_files(namespace: str)이기 때문에 None일 수 있는 값을
    #    직접 넘기지 않고, 빈 문자열로 정규화 후 전달한다.
    #  - 빈 문자열("")은 ingest_local_files 내부 _compute_effective_ns()
    #    에서 falsy로 처리되어 CHROMA_NAMESPACE_LOCAL / "<slug>-local"
    #    규칙으로 다시 보정된다.
    ns_for_ingest: str = namespace or ""
    json_paths, _preview, chunk_total = ingest_local_files(
        globs=patterns,
        namespace=ns_for_ingest,
        persist_directory=persist_directory,
        topic_slug=slug,
        root_dir=str(current_path()),
        add_web_pages_json_to_chroma=_add_to_chroma,
        web_page_json_to_documents=_to_docs,
    )
    logger.info("[LOCAL RAG] findings quick-ingest: paths=%s, chunks=%d", json_paths, chunk_total)
    return int(chunk_total)

# ──────────────────────────────────────────────────────────────
# Quick wrapper for routers: callable-safe, narrow signature
# ──────────────────────────────────────────────────────────────
def quick_ingest_findings(topic_slug: str, patterns: Optional[list[str]] = None) -> int:
    """
    라우터(after_synth 등)에서 findings를 빠르게 인덱싱하기 위한 얇은 래퍼.
    - 주어진 patterns가 없으면 add_local_findings_to_chroma의 기본 패턴을 사용
    - 반환: 인덱싱된 청크 수(int)
    """
    try:
        return add_local_findings_to_chroma(topic_slug=topic_slug, patterns=patterns)
    except TypeError:
        # 혹시 키워드 인자를 지원하지 않는 이전 시그니처가 있을 경우 폴백
        if patterns is not None:
            return add_local_findings_to_chroma(topic_slug, patterns=patterns)
        return add_local_findings_to_chroma(topic_slug)
