# -*- coding: utf-8 -*-
"""§13 fix verification — _ensure_slide_number() 동작 확인 (LLM 호출 0).

합성 spec 으로 4개 layout (0/1/2/5) 슬라이드 생성 → 결과 .pptx 재오픈해서
SLIDE_NUMBER placeholder 가 들어왔는지·SECTION_HEADER 만 의도적 제외인지 검증.
"""
from __future__ import annotations
import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

from pptx import Presentation
from agent.export.renderer import render_deck
from agent.export.spec import SlideDeckSpec, SlideSpec, TableSpec

EMU = 360000
TEMPLATE = Path("templates/agency_default.pptx")
OUT = Path("reports/_cli_test/synthetic_v4.pptx")


def main() -> None:
    spec = SlideDeckSpec(
        slug="synthetic_v4",
        topic_title="SLIDE_NUMBER fix 검증",
        slides=[
            SlideSpec(layout_id=0, title="페이지 번호 fix 검증",
                      body="합성 spec — LLM 호출 0"),
            SlideSpec(layout_id=2, title="1. SECTION_HEADER (의도적 제외)",
                      body="여기엔 페이지 번호 없어야 함"),
            SlideSpec(layout_id=1, title="TITLE_CONTENT bullets",
                      bullets=["bullet 1", "bullet 2", "bullet 3"]),
            SlideSpec(layout_id=1, title="TITLE_TABLE (layout 5 dispatch)",
                      table=TableSpec(header=["A", "B"], rows=[["1", "2"]])),
        ],
    )

    out = render_deck(spec, template_path=TEMPLATE, out_path=OUT)
    print(f"[ok] rendered → {out} ({out.stat().st_size:,} B)")

    # 재오픈 검증
    prs = Presentation(str(out))
    print(f"\n=== 재오픈 검증 (n_slides={len(prs.slides)}) ===")
    expected = {
        0: ("TITLE", True),
        1: ("SECTION_HEADER", False),  # 의도적 제외
        2: ("TITLE_CONTENT", True),
        3: ("TITLE_TABLE", True),
    }
    all_ok = True
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        has_slide_num = False
        slide_num_info = None
        for ph in slide.placeholders:
            try:
                if int(ph.placeholder_format.type) == 13:
                    has_slide_num = True
                    slide_num_info = (
                        f"idx={ph.placeholder_format.idx} "
                        f"top={ph.top/EMU:.2f}cm left={ph.left/EMU:.2f}cm "
                        f"w={ph.width/EMU:.2f}cm h={ph.height/EMU:.2f}cm"
                    )
                    break
            except Exception:
                continue
        exp_layout, exp_has = expected[i]
        layout_ok = layout_name == exp_layout
        has_ok = has_slide_num == exp_has
        status = "OK" if (layout_ok and has_ok) else "FAIL"
        if status == "FAIL":
            all_ok = False
        print(f"  S{i+1} layout={layout_name!r:<20} "
              f"SLIDE_NUMBER={'YES' if has_slide_num else 'NO ':<3} "
              f"expect=({exp_layout}, {'YES' if exp_has else 'NO'})  [{status}]")
        if slide_num_info:
            print(f"       └─ {slide_num_info}")

    print(f"\n결론: {'모든 슬라이드 OK' if all_ok else 'FAIL — 위 항목 확인'}")
    sys.exit(0 if all_ok else 1)


if __name__ == "__main__":
    main()
