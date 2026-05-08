# -*- coding: utf-8 -*-
"""§13 fix - SLIDE_NUMBER placeholder 누락 진단.

목적: add_slide(layout) 직후 slide 인스턴스에 SLIDE_NUMBER placeholder 가
들어오는지, 들어온다면 idx/top 좌표가 무엇인지 확인.
"""
from __future__ import annotations
import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

from pptx import Presentation

EMU = 360000
TEMPLATE = Path("templates/agency_default.pptx")


def cm(emu) -> str:
    if emu is None:
        return "None"
    return f"{emu / EMU:.2f}cm"


def dump_placeholders(label: str, container) -> None:
    print(f"  [{label}] placeholders:")
    found = False
    for ph in container.placeholders:
        found = True
        pf = ph.placeholder_format
        print(
            f"    idx={pf.idx:>2}  type={str(pf.type):<30} "
            f"name={ph.name!r:<35} "
            f"top={cm(ph.top):>8} left={cm(ph.left):>8} "
            f"w={cm(ph.width):>8} h={cm(ph.height):>8}"
        )
    if not found:
        print("    (none)")


def dump_shapes(label: str, container) -> None:
    print(f"  [{label}] shapes (non-placeholder):")
    found = False
    for sh in container.shapes:
        if sh.is_placeholder:
            continue
        found = True
        try:
            top = cm(sh.top)
            left = cm(sh.left)
            w = cm(sh.width)
            h = cm(sh.height)
        except Exception:
            top = left = w = h = "?"
        print(
            f"    name={sh.name!r:<40} type={sh.shape_type} "
            f"top={top} left={left} w={w} h={h}"
        )
    if not found:
        print("    (none)")


def main() -> None:
    if not TEMPLATE.exists():
        print(f"ERROR: template missing: {TEMPLATE}")
        sys.exit(1)

    print(f"=== template: {TEMPLATE} ===\n")
    prs = Presentation(str(TEMPLATE))

    print(f"slide_layouts count: {len(prs.slide_layouts)}\n")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"--- layout[{i}] name={layout.name!r}")
        dump_placeholders(f"layout[{i}]", layout)
        dump_shapes(f"layout[{i}]", layout)
        print()

    # 4개 layout (0/1/2/5) 에 대해 add_slide 후 슬라이드 인스턴스에 어떤 placeholder 가 들어오는지 확인
    print("\n=== add_slide() 직후 슬라이드 인스턴스 placeholder/shape 검사 ===\n")
    for layout_idx in (0, 1, 2, 5):
        if layout_idx >= len(prs.slide_layouts):
            print(f"layout[{layout_idx}] not present, skip")
            continue
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        print(f"--- add_slide(layout[{layout_idx}] name={layout.name!r})")
        dump_placeholders(f"slide-after-add (layout {layout_idx})", slide)
        dump_shapes(f"slide-after-add (layout {layout_idx})", slide)
        print()


if __name__ == "__main__":
    main()
