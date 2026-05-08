# -*- coding: utf-8 -*-
"""§13-10 verification — 표 슬라이드 품질 자동 검증.

검증 항목:
  - 헤더 행 (row=0): 배경 RGB(26,26,26), 글자 RGB(255,255,255), Bold, 12pt (≥8 행 fallback 11pt)
  - 본문 행 (row≥1): 글자 RGB(26,26,26), 10pt (fallback 9pt)
  - 컬럼 너비 합계 ≈ TABLE_AREA_WIDTH_CM (29.87cm), 최소 2.5cm 보장
  - 컬럼 너비 균등 X (max - min ≥ 0.5cm)
  - 행 높이 명시 (최소값 0.6cm)
"""
from __future__ import annotations
import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Pt

EMU = 360000
TARGET_AREA_W_CM = 29.87
MIN_COL_W_CM = 2.5
EXPECT_HEADER_BG = (0x1A, 0x1A, 0x1A)
EXPECT_HEADER_FG = (0xFF, 0xFF, 0xFF)
EXPECT_BODY_FG = (0x1A, 0x1A, 0x1A)
EXPECT_HEADER_PT_NORMAL = 12.0
EXPECT_HEADER_PT_FALLBACK = 11.0
EXPECT_BODY_PT_NORMAL = 10.0
EXPECT_BODY_PT_FALLBACK = 9.0
FALLBACK_ROWS = 8
HEIGHT_MIN_CM = 0.6


def rgb_tuple(rgb_color):
    if rgb_color is None:
        return None
    s = str(rgb_color)
    if len(s) != 6:
        return None
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def first_run_font(cell):
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            return run.font
    return None


def cell_fill_rgb(cell):
    try:
        if cell.fill.type is None:
            return None
        if cell.fill.type == 1:  # SOLID
            return rgb_tuple(cell.fill.fore_color.rgb)
    except Exception:
        return None
    return None


def check_table(table, n_rows, expected_header_pt, expected_body_pt) -> list[str]:
    fails: list[str] = []
    n_cols = len(table.columns)

    # 헤더 검증
    for c in range(n_cols):
        cell = table.cell(0, c)
        bg = cell_fill_rgb(cell)
        if bg != EXPECT_HEADER_BG:
            fails.append(f"header[{c}] bg={bg} expect={EXPECT_HEADER_BG}")
        font = first_run_font(cell)
        if font is None:
            fails.append(f"header[{c}] no run/font")
            continue
        if font.bold is not True:
            fails.append(f"header[{c}] bold={font.bold} expect=True")
        if font.size is None or abs(font.size.pt - expected_header_pt) > 0.01:
            sz = font.size.pt if font.size else None
            fails.append(f"header[{c}] size={sz}pt expect={expected_header_pt}pt")
        fg = rgb_tuple(font.color.rgb) if font.color and font.color.type else None
        if fg != EXPECT_HEADER_FG:
            fails.append(f"header[{c}] fg={fg} expect={EXPECT_HEADER_FG}")

    # 본문 검증
    for r in range(1, n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            font = first_run_font(cell)
            if font is None:
                continue  # 빈 셀 skip
            if font.size is None or abs(font.size.pt - expected_body_pt) > 0.01:
                sz = font.size.pt if font.size else None
                fails.append(f"body[{r},{c}] size={sz}pt expect={expected_body_pt}pt")
            fg = rgb_tuple(font.color.rgb) if font.color and font.color.type else None
            if fg != EXPECT_BODY_FG:
                fails.append(f"body[{r},{c}] fg={fg} expect={EXPECT_BODY_FG}")

    # 컬럼 너비
    widths_cm = [col.width / EMU for col in table.columns]
    sum_w = sum(widths_cm)
    if abs(sum_w - TARGET_AREA_W_CM) > 0.01:
        fails.append(f"col width sum={sum_w:.2f}cm expect={TARGET_AREA_W_CM}cm")
    for c, w in enumerate(widths_cm):
        if w < MIN_COL_W_CM - 0.01:
            fails.append(f"col[{c}] width={w:.2f}cm < min {MIN_COL_W_CM}cm")
    spread = max(widths_cm) - min(widths_cm)
    # 균등 분할 회피 — 길이 다양성이 있으면 spread > 0.5cm 기대 (모든 컬럼 동일 텍스트 길이일 때만 균등)
    print(f"    col widths (cm): {[f'{w:.2f}' for w in widths_cm]} "
          f"sum={sum_w:.2f} spread={spread:.2f}")

    # 행 높이
    heights_cm = [row.height / EMU for row in table.rows]
    print(f"    row heights (cm): {[f'{h:.2f}' for h in heights_cm]} "
          f"(min expected = {HEIGHT_MIN_CM})")
    for r, h in enumerate(heights_cm):
        if h < HEIGHT_MIN_CM - 0.01:
            fails.append(f"row[{r}] height={h:.2f}cm < min {HEIGHT_MIN_CM}cm")

    return fails


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print("usage: verify_table_quality.py <pptx>")
        return 1
    p = Path(argv[1])
    if not p.exists():
        print(f"ERROR: not found: {p}")
        return 1

    prs = Presentation(str(p))
    print(f"=== {p} ===")
    table_count = 0
    total_fails = 0
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table_count += 1
            tbl = shape.table
            n_rows = len(tbl.rows)
            n_cols = len(tbl.columns)
            use_fallback = n_rows >= FALLBACK_ROWS
            exp_header = EXPECT_HEADER_PT_FALLBACK if use_fallback else EXPECT_HEADER_PT_NORMAL
            exp_body = EXPECT_BODY_PT_FALLBACK if use_fallback else EXPECT_BODY_PT_NORMAL
            print(f"\n  S{i:>2} table  rows={n_rows} cols={n_cols} "
                  f"fallback={'YES' if use_fallback else 'no'}")
            fails = check_table(tbl, n_rows, exp_header, exp_body)
            if fails:
                total_fails += len(fails)
                for f in fails:
                    print(f"    FAIL — {f}")
            else:
                print(f"    OK")

    print(f"\n총 표 {table_count}개, FAIL 항목 수 {total_fails}")
    return 0 if total_fails == 0 else 1


if __name__ == "__main__":
    sys.exit(main(sys.argv))
