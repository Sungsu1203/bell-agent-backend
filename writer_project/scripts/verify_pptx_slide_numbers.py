# -*- coding: utf-8 -*-
"""§13 fix verification — 임의 .pptx 의 슬라이드별 페이지 번호 placeholder 유무 검증."""
from __future__ import annotations
import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

from pptx import Presentation

EMU = 360000


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print("usage: verify_pptx_slide_numbers.py <pptx>")
        return 1
    p = Path(argv[1])
    if not p.exists():
        print(f"ERROR: not found: {p}")
        return 1

    prs = Presentation(str(p))
    print(f"=== {p} (n_slides={len(prs.slides)}) ===\n")
    counts = {"with": 0, "without_excl_section": 0, "section_header": 0}
    for i, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name
        slide_num_ph = None
        for ph in slide.placeholders:
            try:
                if int(ph.placeholder_format.type) == 13:
                    slide_num_ph = ph
                    break
            except Exception:
                continue
        has = slide_num_ph is not None
        if layout_name == "SECTION_HEADER":
            counts["section_header"] += 1
            tag = "EXEMPT"
        elif has:
            counts["with"] += 1
            tag = "OK"
        else:
            counts["without_excl_section"] += 1
            tag = "MISSING"
        info = ""
        if slide_num_ph is not None:
            info = (f"  idx={slide_num_ph.placeholder_format.idx} "
                    f"top={slide_num_ph.top/EMU:.2f}cm "
                    f"left={slide_num_ph.left/EMU:.2f}cm")
        print(f"  S{i:>2} layout={layout_name!r:<25} [{tag}]{info}")
    print(f"\n요약: with={counts['with']}  "
          f"without (non-section)={counts['without_excl_section']}  "
          f"SECTION_HEADER (exempt)={counts['section_header']}")
    if counts["without_excl_section"] > 0:
        print("FAIL — SECTION_HEADER 외 슬라이드에 SLIDE_NUMBER 누락")
        return 1
    print("PASS — 모든 일반 슬라이드에 SLIDE_NUMBER 표시, SECTION_HEADER 만 의도적 제외")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
