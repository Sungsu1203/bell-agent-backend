# tools/local_rag.py
from __future__ import annotations
import logging
logger = logging.getLogger(__name__)

import os, re, json, glob, hashlib
from pathlib import Path
from datetime import datetime
from typing import Any, List, Tuple, Optional
from urllib.parse import unquote

from langchain_core.documents import Document

from core.config import CFG
from core.config import reload_config as reload_config  # 런타임 갱신 허용
from core.paths import (
    current_path,
    research_base_dir,
    research_topic_dir,
    research_resources_dir,   # ← 표준 리소스 경로 헬퍼 사용
)

# ──────────────────────────────────────────────────────────────────────────────
# Optional dependencies: 클래스/함수 핸들을 Any로 보관(없으면 None)
# ──────────────────────────────────────────────────────────────────────────────
try:
    from bs4 import BeautifulSoup as _BeautifulSoup
    BeautifulSoup: Optional[Any] = _BeautifulSoup
except Exception:
    BeautifulSoup = None
    logger.debug("BeautifulSoup not available; falling back to regex for HTML parsing.")

try:
    # 프로젝트 표준: PyPDF2
    from PyPDF2 import PdfReader as _PdfReader
    PdfReader: Optional[Any] = _PdfReader
except Exception:
    PdfReader = None
    logger.debug("PyPDF2 not available; PDF extraction disabled.")

try:
    import docx as _docx  # python-docx
    docx: Optional[Any] = _docx
except Exception:
    docx = None
    logger.debug("python-docx not available; .docx extraction disabled.")

# (선택) Unstructured: PPTX/XLSX 등 포맷을 자동 분해
# unstructured: 있으면 함수/클래스 핸들을 보관, 없으면 None
try:
    from unstructured.partition.auto import partition as _unstructured_partition  # noqa: F401
    partition: Optional[Any] = _unstructured_partition
except Exception:
    partition = None
    logger.debug("unstructured.partition.auto not available.")
try:
    from unstructured.documents.elements import Table as _UnstructuredTable  # noqa: F401
    Table: Optional[Any] = _UnstructuredTable
except Exception:
    Table = None
    logger.debug("unstructured.documents.elements.Table not available.")

# (선택) python-pptx: 슬라이드 단위 샘플 추출용
# python-pptx: 클래스 핸들을 Optional[Any]로 유지
try:
    from pptx import Presentation as _pptx_Presentation  # noqa: F401
    _PptxPresentation: Optional[Any] = _pptx_Presentation
except Exception:
    _PptxPresentation = None
    logger.debug("python-pptx not available; PPTX sampling disabled.")

# ──────────────────────────────────────────────────────────────────────────────
# 내부 유틸 (CFG 우선 → ENV 폴백 단일 진입)
# ──────────────────────────────────────────────────────────────────────────────
# ※ 다른 모듈과의 중복 정의를 피하기 위해, 이미 전역에 있으면 재정의하지 않음
if "_cfg_str" not in globals():
    def _cfg_str(name: str, default: str = "") -> str:
        try:
            v = getattr(CFG, name)
            if v is None:
                raise AttributeError
            s = str(v).strip()
            return s if s != "" else default
        except Exception:
            s = (os.getenv(name, "") or "").strip()
            return s if s != "" else default

if "_cfg_int" not in globals():
    def _cfg_int(name: str, default: int = 0) -> int:
        try:
            v = getattr(CFG, name)
            if v is None or str(v).strip() == "":
                raise AttributeError
            return int(str(v).strip())
        except Exception:
            ev = (os.getenv(name, "") or "").strip()
            try:
                return int(ev) if ev != "" else default
            except Exception:
                return default

if "_cfg_float" not in globals():
    def _cfg_float(name: str, default: float = 0.0) -> float:
        try:
            v = getattr(CFG, name)
            if v is None or str(v).strip() == "":
                raise AttributeError
            return float(str(v).strip())
        except Exception:
            ev = (os.getenv(name, "") or "").strip()
            try:
                return float(ev) if ev != "" else default
            except Exception:
                return default

if "_cfg_bool" not in globals():
    def _cfg_bool(name: str, default: bool = False) -> bool:
        try:
            v = getattr(CFG, name)
            if isinstance(v, bool):
                return v
            s = str(v).strip().lower()
            if s in ("1", "true", "yes", "on"):
                return True
            if s in ("0", "false", "no", "off"):
                return False
            return default
        except Exception:
            s = (os.getenv(name, "") or "").strip().lower()
            if s in ("1", "true", "yes", "on"):
                return True
            if s in ("0", "false", "no", "off"):
                return False
            return default

def _truthy_cfg(name: str, default: bool = False) -> bool:
    return _cfg_bool(name, default)

def _env_int(name: str, default: int) -> int:
    # (하위호환 alias) 기존 호출부 유지
    return _cfg_int(name, default)

def _env_float(name: str, default: float) -> float:
    # (하위호환 alias) 기존 호출부 유지
    return _cfg_float(name, default)

def ensure_config_fresh() -> None:
    """런타임에서 최신 .env/CFG를 반영."""
    try:
        reload_config()
        logger.debug("[LOCAL RAG] reload_config() applied.")
    except Exception:
        # 안전 무시
        pass

def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")

def _sha1_file(path: str) -> str:
    h = hashlib.sha1()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()

def _cache_dir() -> Path:
    cache_dir = _cfg_str("LOCAL_RAG_CACHE_DIR", "")
    if not cache_dir:
        cache_dir = str(Path(research_base_dir()) / ".cache")
    p = Path(cache_dir)
    p.mkdir(parents=True, exist_ok=True)
    return p

def _cache_paths(path: str) -> Tuple[str, str]:
    """대용량 변환 캐시 경로 결정: cache/{sha1}.{ext}.json"""
    cdir = _cache_dir()
    sha1 = _sha1_file(path)
    ext = Path(path).suffix.lower().lstrip(".")
    meta = (cdir / f"{sha1}.{ext}.meta.json").as_posix()
    data = (cdir / f"{sha1}.{ext}.data.json").as_posix()
    return meta, data

def _cache_load(path: str) -> Optional[List[dict]]:
    # CFG 우선 → ENV 폴백 (신규 일원화 헬퍼 사용)
    if _truthy_cfg("LOCAL_RAG_CACHE_IGNORE", False):
        return None
    meta, data = _cache_paths(path)
    try:
        if os.path.isfile(meta) and os.path.isfile(data):
            with open(meta, "r", encoding="utf-8") as fm:
                _ = json.load(fm)  # 향후 버전 검증 등
            with open(data, "r", encoding="utf-8") as fd:
                items = json.load(fd)
            logger.info("[LOCAL RAG] cache hit → %s", data)
            return items
    except Exception as e:
        logger.debug("cache load failed: %s", e)
    return None

def _cache_save(path: str, items_from_reader: List[dict]) -> None:
    try:
        meta, data = _cache_paths(path)
        with open(meta, "w", encoding="utf-8") as fm:
            json.dump({
                "path": str(Path(path).resolve()),
                "sha1": _sha1_file(path),
                "saved_at": _now_iso(),
                "count": len(items_from_reader),
            }, fm, ensure_ascii=False, indent=2)
        with open(data, "w", encoding="utf-8") as fd:
            json.dump(items_from_reader, fd, ensure_ascii=False, indent=2)
        logger.info("[LOCAL RAG] cache saved → %s", data)
    except Exception as e:
        logger.debug("cache save failed: %s", e)

def _file_version(path: str) -> str:
    """파일 변경을 반영하는 버전 식별자."""
    forced = _cfg_str("LOCAL_RAG_FORCE_VERSION", "")
    if forced:
        logger.debug("Using forced version for %s: %s", path, forced)
        return str(forced)

    mode = _cfg_str("LOCAL_RAG_VERSION_MODE", "mtime").lower()
    try:
        if mode == "sha1":
            h = hashlib.sha1()
            with open(path, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    h.update(chunk)
            ver = h.hexdigest()[:12]
        else:
            st = os.stat(path)
            ver = f"{int(st.st_mtime)}_{st.st_size}"
        logger.debug("Computed version for %s: %s (mode=%s)", path, ver, mode)
        return ver
    except Exception as e:
        logger.warning("Version compute failed for %s: %s", path, e)
        return "na"

def _truncate(s: str, max_chars_env: str = "LOCAL_RAG_MAX_TEXT_CHARS") -> str:
    # CFG 우선 → ENV
    cfg_val = getattr(CFG, max_chars_env, None)
    if cfg_val is None or str(cfg_val).strip() == "":
        max_chars = _cfg_int(max_chars_env, 200000)
    else:
        max_chars = int(str(cfg_val))
    if max_chars > 0 and len(s) > max_chars:
        return s[:max_chars]
    return s

# ── Readers ──────────────────────────────────────────────────────────────────
def _read_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _read_md(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _read_csv(path: str) -> str:
    # 간단 CSV → 줄 단위 텍스트(필요시 unstructured/판다스로 확장)
    try:
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return Path(path).read_text(encoding="cp949", errors="ignore")

def _read_html(path: str) -> str:
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    if BeautifulSoup is not None:
        try:
            return BeautifulSoup(raw, "lxml").get_text(" ", strip=True)
        except Exception:
            return BeautifulSoup(raw, "html.parser").get_text(" ", strip=True)
    txt = re.sub(r"<[^>]+>", " ", raw)
    return re.sub(r"\s+", " ", txt).strip()

def _read_docx(path: str) -> str:
    if not docx:
        raise RuntimeError("python-docx 미설치")
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs)

def _read_pdf_pages(path: str, max_pages: Optional[int] = None) -> List[str]:
    """
    PDF 페이지 텍스트를 페이지별 리스트로 반환.
    max_pages 지정 없으면 ENV/CFG의 LOCAL_RAG_PDF_MAX_PAGES(기본 30) 사용.
    """
    if not PdfReader:
        raise RuntimeError("PyPDF2 미설치")
    reader = PdfReader(path)
    limit = int(max_pages) if (isinstance(max_pages, int) and max_pages > 0) else _env_int("LOCAL_RAG_PDF_MAX_PAGES", 30)
    pages: List[str] = []
    for i, p in enumerate(reader.pages, start=1):
        if limit and i > limit:
            break
        try:
            txt = p.extract_text() or ""
            pages.append(txt)
        except Exception:
            pages.append("")
    return pages

# ── PPTX helpers ─────────────────────────────────────────────────────────────
def _read_unstructured_elements(path: str) -> List[dict]:
    """
    PPTX, XLSX 등을 unstructured로 분해하여 요소(Element) 리스트로 반환.
    각 요소 dict 예시: {"sheet": "<슬라이드/시트/파트>", "row_index": <번호>, "content": "<텍스트>"}.
    """
    if partition is None:
        raise RuntimeError("unstructured 라이브러리 미설치")
    try:
        elements = partition(filename=path)
    except Exception as e:
        logger.warning("Unstructured partition failed for %s: %s", path, e)
        return []

    items: List[dict] = []
    for i, element in enumerate(elements):
        content = str(getattr(element, "text", "") or "").strip()
        if Table is not None and isinstance(element, Table):
            text_as_csv = getattr(element, "text_as_csv", "") or ""
            if text_as_csv:
                content = (content + "\n\n" + text_as_csv).strip()
        if not content:
            continue
        metadata = getattr(element, "metadata", None)
        part_name = getattr(metadata, "page_number", None) or getattr(metadata, "sheet_name", None) or "Part"
        items.append({"sheet": str(part_name), "row_index": i + 1, "content": content})
    return items

def _pick_sample_indices(n: int, ratio: float) -> List[int]:
    import math
    k = max(1, int(math.ceil(n * max(0.01, min(1.0, ratio)))))
    base = {0, max(0, n - 1)}
    if n <= len(base):
        return sorted(base)
    step = max(1, n // k)
    for i in range(0, n, step):
        base.add(i)
        if len(base) >= k:
            break
    return sorted(min(idx, n - 1) for idx in base)

def _pptx_extract_titles_bullets(path: str, sample_large: bool = False) -> List[dict]:
    """
    대용량 PPTX 전용: 제목+불릿만 추출, 표/차트 존재는 태그로 표기.
    sample_large=True면 슬라이드 일부 샘플링(기본 20%), 항상 1/마지막 포함.
    """
    if _PptxPresentation is None:
        raise RuntimeError("python-pptx 미설치")

    prs = _PptxPresentation(path)
    slides = list(prs.slides)
    total = len(slides)
    ratio = _env_float("LOCAL_RAG_SAMPLE_RATIO", 0.2) if sample_large else 1.0
    pick = _pick_sample_indices(total, ratio)

    items: List[dict] = []
    for sidx in pick:
        slide = slides[sidx]
        texts: List[str] = []

        # 제목 추출
        title = ""
        try:
            if getattr(slide.shapes, "title", None) is not None:
                title = str(getattr(slide.shapes.title, "text", "") or "").strip()
        except Exception:
            title = ""

        if title:
            texts.append(title)

        # 불릿 추출
        try:
            for shp in slide.shapes:
                has_tf = bool(getattr(shp, "has_text_frame", False))
                tf = getattr(shp, "text_frame", None)
                if not (has_tf and tf):
                    continue
                for p in getattr(tf, "paragraphs", []) or []:
                    raw = str(getattr(p, "text", "") or "").strip()
                    if not raw:
                        continue
                    if title and raw == title:
                        continue
                    level = int(getattr(p, "level", 0) or 0)
                    if level > 0 or len(raw) >= 2:
                        texts.append(f"- {raw}")
        except Exception:
            pass

        # 표/차트 존재 태그
        try:
            has_table = any(bool(getattr(shp, "has_table", False)) for shp in slide.shapes)
            has_chart = any(getattr(shp, "chart", None) is not None for shp in slide.shapes)
            if has_table:
                texts.append("[표 있음]")
            if has_chart:
                texts.append("[차트 있음]")
        except Exception:
            pass

        content = "\n".join(t for t in texts if t).strip()
        if not content:
            continue
        items.append({"sheet": str(sidx + 1), "row_index": sidx + 1, "content": content})

    logger.info("[LOCAL RAG] PPTX %s → slim extracted (slides=%d/%d)", Path(path).name, len(items), total)
    return items

def _read_pptx(path: str) -> List[dict]:
    """
    우선 python-pptx로 슬라이드 단위(제목+불릿) 추출 시도,
    실패/미설치 시 unstructured로 폴백.
    """
    if _truthy_cfg("SKIP_PPTX", False):
        logger.info("SKIP_PPTX=1 → _read_pptx skipped for %s", path)
        return []

    if _PptxPresentation is not None:
        try:
            return _pptx_extract_titles_bullets(path, sample_large=False)
        except Exception as e:
            logger.warning("python-pptx extract failed, fallback to unstructured: %s", e)

    # fallback
    return _read_unstructured_elements(path)

# ── Markdown-like chunking ────────────────────────────────────────────────────
_MD_FENCE_RE = re.compile(r"```.*?```", flags=re.DOTALL)
_MD_TABLE_LINE_RE = re.compile(r"^\s*\|.*\|\s*$")

def _summarize_table_lines(lines: List[str], max_rows: int = 3) -> str:
    head = [ln.strip() for ln in lines[:max_rows] if ln.strip()]
    if not head:
        return "[표 요약]"
    sample = " / ".join(head)
    return f"[표 요약: {sample[:120]}...]"

def _split_markdown_to_chunks(
    text: str,
    min_chars: Optional[int] = None,
    max_chars: Optional[int] = None,
) -> List[str]:
    """
    규칙:
      - ```fence``` 코드블록은 '[코드 요약]'으로 치환
      - 파이프(|) 기반 테이블은 '[표 요약: ...]' 1줄로 치환
      - 빈줄 기준 문단 → 600~1,200자 범위로 병합/분할(CFG/ENV로 조정)
    """
    min_chars = min_chars or _cfg_int("LOCAL_RAG_MIN_CHARS", 600)
    max_chars = max_chars or _cfg_int("LOCAL_RAG_MAX_CHARS", 1200)

    if _cfg_str("LOCAL_RAG_CHUNK_MODE", "paragraph").lower() == "simple":
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        return lines

    # 1) 코드블록 요약 치환
    text = _MD_FENCE_RE.sub("[코드 요약]", text)

    # 2) 테이블 블록을 한 덩어리로 요약 치환
    lines = text.splitlines()
    out_lines: List[str] = []
    buf_table: List[str] = []
    def _flush_table():
        nonlocal buf_table, out_lines
        if buf_table:
            out_lines.append(_summarize_table_lines(buf_table))
            buf_table = []
    for ln in lines:
        if _MD_TABLE_LINE_RE.match(ln):
            buf_table.append(ln)
        else:
            _flush_table()
            out_lines.append(ln)
    _flush_table()
    text = "\n".join(out_lines)

    # 3) 빈줄 기준 문단
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    # 4) 병합(짧은 문단 → 다음 문단과 붙임)
    merged, buf = [], ""
    for p in paras:
        if len(buf) + len(p) + 1 < min_chars:
            buf = (buf + " " + p).strip() if buf else p
        else:
            if buf:
                merged.append(buf)
            buf = p
    if buf:
        merged.append(buf)

    # 5) 최대 길이 초과는 단어 경계로 슬라이스
    chunks: List[str] = []
    for m in merged:
        cur = m.strip()
        while len(cur) > max_chars:
            cut = cur[:max_chars]
            cut_idx = cut.rfind(" ")
            if cut_idx < max_chars * 0.5:  # 단어 경계가 너무 앞이면 하드컷
                cut_idx = max_chars
            chunks.append(cur[:cut_idx].strip())
            cur = cur[cut_idx:].strip()
        if cur:
            chunks.append(cur)
    return chunks

def _file_uri(path: str) -> str:
    return Path(path).resolve().as_uri()  # file:/// 형태 보장

# ──────────────────────────────────────────────────────────────────────────────
# 변환: 로컬 파일 → web.json 아이템 배열
# ──────────────────────────────────────────────────────────────────────────────
def _to_webjson_items(path: str, *, max_pages_per_file: Optional[int] = None) -> List[dict]:
    ext = Path(path).suffix.lower()
    title = Path(path).name
    url_click = _file_uri(path)  # 사람이 눌러 열어볼 주소
    ver = _file_version(path)
    fetched_at = _now_iso()

    _ct_map = {
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".html": "text/html",
        ".htm": "text/html",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv": "text/csv",
    }
    content_type = _ct_map.get(ext, "application/octet-stream")

    # 용량 가드(+ 대용량 PPTX 슬림/샘플링)
    max_mb = _cfg_float("LOCAL_RAG_MAX_FILE_MB", 50.0)
    try:
        size_mb = os.path.getsize(path) / (1024 * 1024)
    except Exception:
        size_mb = None
    is_large = bool(size_mb and max_mb > 0 and size_mb > max_mb)
    enable_sample = _truthy_cfg("LOCAL_RAG_SAMPLE_LARGE", False)

    text: Optional[str] = None
    items_from_reader: Optional[List[dict]] = None

    try:
        if ext == ".pdf":
            if is_large:
                logger.info("[LOCAL RAG] skip large PDF (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            pages = _read_pdf_pages(path, max_pages=max_pages_per_file)
            items_from_reader = []
            for i, txt in enumerate(pages, start=1):
                if (txt or "").strip():
                    items_from_reader.append({"page_num": i, "content": txt})

        elif ext in (".txt", ".md", ".markdown", ".html", ".htm", ".docx", ".csv"):
            if is_large:
                logger.info("[LOCAL RAG] skip large text-like (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            if ext == ".txt":
                text = _read_txt(path)
            elif ext in (".md", ".markdown"):
                text = _read_md(path)
            elif ext in (".html", ".htm"):
                text = _read_html(path)
            elif ext == ".docx":
                text = _read_docx(path)
            elif ext == ".csv":
                text = _read_csv(path)

        elif ext == ".pptx":
            if is_large and enable_sample:
                cached = _cache_load(path)
                if cached is not None:
                    items_from_reader = cached
                    logger.info("[LOCAL RAG] sampled/slim extracted (cache) → %s", Path(path).name)
                else:
                    items_from_reader = _pptx_extract_titles_bullets(path, sample_large=True)
                    _cache_save(path, items_from_reader or [])
                    logger.info("[LOCAL RAG] sampled/slim extracted (fresh) → %s", Path(path).name)
            elif is_large and not enable_sample:
                logger.info("[LOCAL RAG] skip large file (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            else:
                prefer_slim = _truthy_cfg("LOCAL_RAG_PPTX_SLIM", False)
                if prefer_slim and _PptxPresentation is not None:
                    items_from_reader = _pptx_extract_titles_bullets(path, sample_large=False)
                else:
                    items_from_reader = _read_unstructured_elements(path)

        elif ext == ".xlsx":
            if is_large:
                logger.info("[LOCAL RAG] skip large XLSX (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            items_from_reader = _read_unstructured_elements(path)

        else:
            logger.debug("[LOCAL RAG] unsupported extension skipped: %s", path)
            return []

    except RuntimeError as e:
        logger.warning("[LOCAL RAG] Reader dependency failed (RuntimeError): %s -> %s", path, e)
        return []
    except Exception as e:
        logger.warning("[LOCAL RAG] Extraction error: %s -> %s", path, e)

    # 2) 결과 통합 및 메타데이터 생성 (공통)
    final_items: List[dict] = []

    # 리스트 형태 결과(PDF, PPTX, XLSX)
    if items_from_reader is not None and isinstance(items_from_reader, list):
        minc = _cfg_int("LOCAL_RAG_MIN_CHARS", 600)
        maxc = _cfg_int("LOCAL_RAG_MAX_CHARS", 1200)
        for it in items_from_reader:
            part_label = str(it.pop("sheet", it.pop("page_num", "Part")))
            index_num = it.pop("row_index", it.pop("page_num", 1))
            content = (it.pop("content", "") or "").strip()
            if not content:
                continue
            chunks = _split_markdown_to_chunks(content, minc, maxc)
            if not chunks:
                continue
            for j, ch in enumerate(chunks, start=1):
                ch = _truncate(ch, "LOCAL_RAG_MAX_TEXT_CHARS")
                final_items.append({
                    "title": f"{title} ({part_label}, Index: {index_num}, Chunk {j})",
                    "url": f"{url_click}#part={part_label}&index={index_num}",
                    "source": f"{url_click}__s_{part_label}__r_{index_num}__c_{j}__v_{ver}",
                    "content": ch,
                    "content_type": content_type,
                    "fetched_at": fetched_at,
                })

    # 단일 텍스트 결과(TXT, MD, HTML, DOCX, CSV)
    elif text is not None:
        content = _truncate(text or "", "LOCAL_RAG_MAX_TEXT_CHARS")
        if content.strip():
            minc = _cfg_int("LOCAL_RAG_MIN_CHARS", 600)
            maxc = _cfg_int("LOCAL_RAG_MAX_CHARS", 1200)
            chunks = _split_markdown_to_chunks(content, minc, maxc)
            for j, ch in enumerate(chunks, start=1):
                final_items.append({
                    "title": f"{title} (Chunk {j})",
                    "url": url_click,
                    "source": f"{url_click}__c_{j}__v_{ver}",
                    "content": ch,
                    "content_type": content_type,
                    "fetched_at": fetched_at,
                })

    if not final_items:
        logger.debug("[LOCAL RAG] File yielded no extractable content: %s", path)
        return []

    return final_items

# ──────────────────────────────────────────────────────────────────────────────
# 엔트리: globs → web.json 생성
# ──────────────────────────────────────────────────────────────────────────────
def build_webjson_from_local(
    globs: List[str],
    out_dir: str,
    *,
    max_docs: Optional[int] = None,
    max_pages_per_file: Optional[int] = None,
) -> str:
    """
    globs에 매칭되는 로컬 파일들을 읽어 web.json으로 변환.
    - max_docs: 전체 아이템(청크) 상한. None/0이면 무제한. ENV: LOCAL_RAG_MAX_DOCS
    - max_pages_per_file: 파일별 페이지 상한(PDF 등). None/0이면 무제한. ENV: LOCAL_RAG_MAX_PAGES_PER_FILE
    """
    # 런타임 설정 갱신
    ensure_config_fresh()
    # CFG 우선 → ENV 폴백
    if max_docs is None:
        max_docs = _cfg_int("LOCAL_RAG_MAX_DOCS", 0)
    if max_pages_per_file is None:
        max_pages_per_file = _cfg_int("LOCAL_RAG_MAX_PAGES_PER_FILE", 0)

    logger.info("[LOCAL RAG] CWD: %s", os.getcwd())
    logger.info("[LOCAL RAG] Received globs: %s", globs)
    if max_docs:
        logger.info("[LOCAL RAG] cap: max_docs=%d", max_docs)
    if max_pages_per_file:
        logger.info("[LOCAL RAG] cap: max_pages_per_file=%d", max_pages_per_file)

    Path(out_dir).mkdir(parents=True, exist_ok=True)
    files: List[str] = []

    for g in globs or []:
        g = os.path.expandvars(os.path.expanduser(g))
        logger.debug("[LOCAL RAG] Expanded glob pattern: %s", g)
        matched = glob.glob(g, recursive=True)
        files.extend(matched)
        logger.info("[LOCAL RAG] Pattern %s matched %d files.", g, len(matched))

    files = sorted({f for f in files if os.path.isfile(f)})
    logger.info("[LOCAL RAG] Total unique files found: %d", len(files))

    # [ADD] 우선순위 샘플링: findings.md > pdf > xlsx > pptx > 기타
    def _priority_key(p: str) -> tuple[int, int]:
        name = (p or "").lower()
        # 최우선: 연구 요약/결과물 성격의 findings.md
        if name.endswith("findings.md"):
            return (0, len(name))
        # 그 다음: 포맷별 우선순위
        if name.endswith(".pdf"):
            return (1, len(name))
        if name.endswith(".xlsx"):
            return (2, len(name))
        if name.endswith(".pptx"):
            return (3, len(name))
        return (4, len(name))  # 기타

    files.sort(key=_priority_key)

    # [ADD] cap 적용: CFG.LOCAL_MAX_FILES → ENV LOCAL_MAX_FILES → 기본 1500
    def _cap_int(name: str, default: int) -> int:
        return _cfg_int(name, default)

    _cap = _cap_int("LOCAL_MAX_FILES", 1500)
    if _cap > 0 and len(files) > _cap:
        logger.info("[LOCAL RAG] applying file cap: %d → %d", len(files), _cap)
        files = files[:_cap]


    items: List[dict] = []
    processed_items = 0
    written_items = 0
    total_files = len(files)

    for fi, f in enumerate(files, start=1):
        try:
            file_items = _to_webjson_items(f, max_pages_per_file=max_pages_per_file or None)
        except Exception as e:
            logger.warning("[LOCAL RAG] local ingest 실패: %s -> %s", f, e)
            file_items = []

        for it in file_items:
            # 빈 content 제거
            if not (it.get("content") or "").strip():
                continue
            items.append(it)
            processed_items += 1
            written_items += 1

            # 500개 단위 진행률 로그
            if processed_items % 500 == 0:
                logger.info("[LOCAL RAG] processed %d items (files %d/%d)", processed_items, fi, total_files)

            # 상한 도달 시 조기 종료
            if max_docs and written_items >= max_docs:
                logger.info("[LOCAL RAG] max_docs reached (%d) — truncating", max_docs)
                break

        if max_docs and written_items >= max_docs:
            break

    # 중복 source 제거(최초 등장 우선)
    dedup_seen, dedup_items = set(), []
    for it in items:
        sid = it.get("source")
        if not sid or sid in dedup_seen:
            continue
        dedup_seen.add(sid)
        dedup_items.append(it)
    items = dedup_items

    # 디버그 샘플
    uniq_sources = {it.get("source") for it in items}
    logger.info("[LOCAL RAG] files=%d items=%d unique_sources=%d", len(files), len(items), len(uniq_sources))
    if items:
        def _pretty_src(src: str) -> str:
            s = unquote(src or "")
            if "__v_" in s:
                s = s.split("__v_")[0]
            if "__p_" in s:
                s = s.split("__p_")[0]
            return s

        sample = items[:3]
        sample_titles = [it.get("title", "") for it in sample]
        sample_sources = [_pretty_src(it.get("source", "")) for it in sample]
        sample_urls = [unquote(it.get("url", "")) for it in sample]

        logger.debug("[LOCAL RAG] sample titles : %s", sample_titles)
        logger.debug("[LOCAL RAG] sample sources: %s", sample_sources)
        logger.debug("[LOCAL RAG] sample urls   : %s", sample_urls)

    ts = datetime.now().strftime("%Y_%m_%d_%H%M%S")
    out_path = os.path.join(out_dir, f"local_{ts}.json")
    with open(out_path, "w", encoding="utf-8") as fp:
        json.dump(items, fp, ensure_ascii=False, indent=2)
    logger.info("[LOCAL RAG] web.json saved → %s", out_path)
    return out_path

# ──────────────────────────────────────────────────────────────────────────────
# 파이프: web.json → Chroma 적재 + 미리보기
# ──────────────────────────────────────────────────────────────────────────────
def ingest_local_files(
    globs: List[str],
    namespace: str,
    persist_directory: str | None,
    topic_slug: str,
    root_dir: str,
    add_web_pages_json_to_chroma=None,
    web_page_json_to_documents=None,
) -> Tuple[List[str], List[Document], int]:
    """
    반환: (생성 JSON 경로들, 미리보기 Documents, 인덱싱된 청크 수 합)
    """
    # 런타임 설정 갱신
    ensure_config_fresh()
    # 연구 요약(findings) 포함 옵션
    include_findings = _truthy_cfg("INCLUDE_FINDINGS_IN_VECTOR", False)
    if include_findings:
        slug = (topic_slug or os.getenv("TOPIC_SLUG") or "default").strip()
        findings_dir = research_topic_dir(slug)
        findings_pattern = str(findings_dir / "round-*-findings.md")
        globs = list(globs or [])
        globs.append(findings_pattern)
        logger.info("[LOCAL RAG] findings included → %s", findings_pattern)

        # [ADD] 추가적인 findings 패턴(선택)
        extra_findings = [
            str(findings_dir / "*findings*.md"),
            str(findings_dir / "findings.md"),
        ]
        for pat in extra_findings:
            if pat not in globs:
                globs.append(pat)
        logger.debug("[LOCAL RAG] extra findings patterns included: %s", extra_findings)

    if not globs:
        logger.info("[LOCAL RAG] no globs provided → skip ingest")
        return ([], [], 0)

    # 표준 리소스 디렉터리 정책 사용
    res_dir = str(research_resources_dir(topic_slug or _cfg_str("TOPIC_SLUG", "default")))
    json_path = build_webjson_from_local(globs, res_dir)

    chunk_total = 0
    if add_web_pages_json_to_chroma is not None:
        try:
            _orig, chunk_count = add_web_pages_json_to_chroma(
                json_path, namespace=namespace, persist_directory=persist_directory
            )
            chunk_total += int(chunk_count or 0)
            logger.info("[LOCAL RAG] added to chroma: chunks=%s (ns=%s, dir=%s)", chunk_count, namespace, persist_directory)
        except Exception as e:
            logger.warning("[LOCAL RAG] add_web_pages_json_to_chroma(local) 실패: %s", e)

    docs_preview: List[Document] = []
    if web_page_json_to_documents is not None:
        try:
            docs_preview = web_page_json_to_documents(json_path)[:8]
            logger.debug("[LOCAL RAG] preview docs: %d", len(docs_preview))
        except Exception as e:
            logger.warning("[LOCAL RAG] preview build(local) 실패: %s", e)

    return ([json_path], docs_preview, chunk_total)
