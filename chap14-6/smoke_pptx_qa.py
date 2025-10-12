# -*- coding: utf-8 -*-
from __future__ import annotations
import os, sys
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from langchain_core.documents import Document

import logging
logger = logging.getLogger(__name__)

import re
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser

from core.llm import get_llm


# 프로젝트 내 도구 재사용
from tools.web_rag import (
    documents_to_chroma,
    retrieve,
    clear_vector_store,
    _default_chroma_dir,
)

def get_pptx_path(argv: Optional[List[str]] = None, env_key: str = "PPTX_PATH") -> Path:
    argv = argv or sys.argv
    arg: Optional[str] = argv[1] if len(argv) >= 2 and argv[1].strip() else None
    env: Optional[str] = os.getenv(env_key)
    source = arg or (env.strip() if env else None)

    if source is None:
        raise SystemExit(
            "PPTX 경로가 없습니다. 사용법: python script.py <pptx_path>  또는  환경변수 PPTX_PATH를 설정하세요."
        )

    return Path(source).expanduser().resolve()

def _shape_texts(shape) -> List[str]:
    out = []
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            txt = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
            if txt.strip():
                out.append(txt.strip())
        elif getattr(shape, "has_table", False) and shape.table:
            tbl = shape.table
            for r in tbl.rows:
                for c in r.cells:
                    if c.text and c.text.strip():
                        out.append(c.text.strip())
        # 그룹 내부 재귀
        if hasattr(shape, "shapes"):
            for s in shape.shapes:
                out.extend(_shape_texts(s))
    except Exception:
        pass
    return out

def load_pptx_to_documents(pptx_path: str) -> List[Document]:
    prs = Presentation(pptx_path)
    docs: List[Document] = []
    p = Path(pptx_path).resolve()
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shp in slide.shapes:
            texts.extend(_shape_texts(shp))
        text = "\n".join(t for t in texts if t).strip()
        if not text:
            continue
        docs.append(Document(
            page_content=text,
            metadata={
                "title": p.name,
                "source": f"file://{p.as_posix()}#slide={i}",
                "slide": i,
            }
        ))
    return docs

def main():
    if len(sys.argv) < 2 and not os.getenv("PPTX_PATH"):
        print("Usage: python smoke_pptx_qa.py <path_to_pptx>  (or set PPTX_PATH)")
        sys.exit(2)

    if not os.getenv("USER_AGENT"):
        os.environ["USER_AGENT"] = "gpt-agent-2025/0.1 (+mailto:you@example.com)"


    # pptx_path = Path(sys.argv[1] if len(sys.argv) >= 2 else os.getenv("PPTX_PATH")).expanduser()
    llm = get_llm()
    pptx_path = get_pptx_path()
    
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(2)

    ns = os.getenv("CHROMA_NAMESPACE", "pptx-test")
    pd = os.getenv("CHROMA_DIR") or _default_chroma_dir(ns)

    # 깨끗한 공간에서 시작
    clear_vector_store(namespace=ns, persist_directory=pd)

    docs = load_pptx_to_documents(str(pptx_path))
    n_src, n_chunks = documents_to_chroma(
        docs,
        namespace=ns,
        persist_directory=pd,
    )
    print(f"loaded: {n_src} docs, {n_chunks} chunks")
    print(f"persist dir = {pd}")

    # 간단 요약(슬라이드 첫 줄들)
    titles = []
    for d in docs:
        head = (d.page_content.splitlines() or [""])[0].strip()
        if head:
            titles.append(head)
    print("\n# Quick summary (slide headings, up to 10)")
    for t in titles[:10]:
        print(f"- {t}")

    # 질의 REPL
    print("\nType your question (enter to quit).")
    while True:
        q = input("Q> ").strip()
        if not q:
            break
        hits = retrieve.invoke({
            "query": q,
            "top_k": 5,
            "namespace": ns,
            "persist_directory": pd,
        })
        print(f"#hits={len(hits)}")
        for d in hits:
            slide = (d.metadata or {}).get("slide")
            snip = (d.page_content or "").replace("\n", " / ")
            print(f"- slide {slide}: {snip[:160]}...")

        try:
            ans = _synthesize_answer(q, hits,llm)
            print("\nAnswer:\n" + ans)
        except Exception as e:
            print(f"\n[WARN] synthesis failed: {e}")

# def _synthesize_answer(q, hits, k=5):

def _synthesize_answer(q: str, hits, llm, k: int = 5) -> str:
    """
    검색 상위 hits로 간단 답변 생성. 항상 문자열을 반환한다.
    """
    if not hits:
        return "(no context)"

    # 1) 슬라이드 번호 우선
    m = re.search(r"\bslide[:\s]*0*(\d+)\b", q, re.I)
    slide_hits = hits
    if m:
        want = int(m.group(1))
        slide_hits = [d for d in hits if (getattr(d, "metadata", {}) or {}).get("slide") == want] + \
                     [d for d in hits if (getattr(d, "metadata", {}) or {}).get("slide") != want]

    # 2) 컨텍스트
    ctx = "\n\n---\n\n".join(
        f"(slide {(getattr(d, 'metadata', {}) or {}).get('slide', '?')})\n{getattr(d, 'page_content', '')}"
        for d in slide_hits[:k]
    )

    # 3) 스타일
    if re.search(r"(한\s*줄|one\s*line)", q, re.I):
        style = "한 문장(최대 25단어)으로 요약하세요."
    elif re.search(r"(표|수치|numbers|figures)", q, re.I):
        style = ("표 수치를 정확히 추출해 '연도 - EV stock(M) / Retired(k) / Recycled(k)' "
                 "형태의 불릿 리스트로 요약하세요. 표에 없는 값은 추측하지 마세요.")
    else:
        style = "핵심만 3줄 이내로 요약하세요."

    # 4) LLM 체인 (항상 str로 변환)
    tpl = "아래 컨텍스트만 사용해 답하세요. 추측 금지.\n{style}\n\n질문: {q}\n\n컨텍스트:\n{ctx}"
    try:
        chain = PromptTemplate.from_template(tpl) | llm | StrOutputParser()
        result = chain.invoke({"q": q, "ctx": ctx, "style": style})
        text = str(result or "").strip()
        return text or "(empty)"
    except Exception as e:
        # 예외가 나도 문자열 반환
        return f"(synthesis failed: {e})"

if __name__ == "__main__":
    main()
