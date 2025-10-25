# tools/local_rag.py
from __future__ import annotations
import logging
logger = logging.getLogger(__name__)

import os, re, json, glob, hashlib
from pathlib import Path
from datetime import datetime
from typing import Any, List, Tuple, Optional
from urllib.parse import unquote

from langchain_core.documents import Document

# ──────────────────────────────────────────────────────────────────────────────
# Optional dependencies: 클래스/함수 핸들을 Any로 보관(없으면 None)
# ──────────────────────────────────────────────────────────────────────────────
try:
    from bs4 import BeautifulSoup as _BeautifulSoup
    BeautifulSoup: Any = _BeautifulSoup
except Exception:
    BeautifulSoup = None
    logger.debug("BeautifulSoup not available; falling back to regex for HTML parsing.")

try:
    # 프로젝트 표준: PyPDF2
    from PyPDF2 import PdfReader as _PdfReader
    PdfReader: Any = _PdfReader
except Exception:
    PdfReader = None
    logger.debug("PyPDF2 not available; PDF extraction disabled.")

try:
    import docx as _docx  # python-docx
    docx: Any = _docx
except Exception:
    docx = None
    logger.debug("python-docx not available; .docx extraction disabled.")

# (선택) Unstructured: PPTX/XLSX 등 포맷을 자동 분해
try:
    from pptx import Presentation as _PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE
except Exception:
    _PptxPresentation = None
    _MSO_SHAPE_TYPE = None
    logger.debug("python-pptx not available; falling back to unstructured for .pptx.")

try:
    from unstructured.partition.auto import partition
    from unstructured.documents.elements import Table  # type: ignore
except Exception:
    partition = None
    Table = None  # type: ignore
    logger.debug("unstructured library not available.")

# (선택) python-pptx: 슬라이드 단위 샘플 추출용
try:
    from pptx import Presentation as _PptxPresentation  # type: ignore[reportMissingTypeStubs]
except Exception:
    _PptxPresentation = None
    logger.debug("python-pptx not available; PPTX sampling disabled.")

# ──────────────────────────────────────────────────────────────────────────────
# 내부 유틸
# ──────────────────────────────────────────────────────────────────────────────
def _env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)))
    except Exception:
        return default

def _env_float(name: str, default: float) -> float:
    try:
        return float(os.getenv(name, str(default)))
    except Exception:
        return default

def _sha1_file(path: str) -> str:
    h = hashlib.sha1()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()

def _cache_paths(path: str) -> Tuple[str, str]:
    """대용량 변환 캐시 경로 결정: cache/{sha1}.{ext}.json"""
    cache_dir = os.getenv("LOCAL_RAG_CACHE_DIR", "cache")
    Path(cache_dir).mkdir(parents=True, exist_ok=True)
    sha1 = _sha1_file(path)
    ext = Path(path).suffix.lower().lstrip(".")
    meta = os.path.join(cache_dir, f"{sha1}.{ext}.meta.json")
    data = os.path.join(cache_dir, f"{sha1}.{ext}.data.json")
    return meta, data

def _cache_load(path: str) -> Optional[List[dict]]:
    if _truthy_env("LOCAL_RAG_CACHE_IGNORE"):
        return None
    meta, data = _cache_paths(path)
    try:
        if os.path.isfile(meta) and os.path.isfile(data):
            with open(meta, "r", encoding="utf-8") as fm:
                _ = json.load(fm)  # 향후 메타 검증용(버전 등)
            with open(data, "r", encoding="utf-8") as fd:
                items = json.load(fd)
            logger.info("[LOCAL RAG] cache hit → %s", data)
            return items  # items_from_reader 포맷
    except Exception as e:
        logger.debug("cache load failed: %s", e)
    return None

def _cache_save(path: str, items_from_reader: List[dict]) -> None:
    try:
        meta, data = _cache_paths(path)
        with open(meta, "w", encoding="utf-8") as fm:
            json.dump({
                "path": str(Path(path).resolve()),
                "sha1": _sha1_file(path),
                "saved_at": _now_iso(),
                "count": len(items_from_reader),
            }, fm, ensure_ascii=False, indent=2)
        with open(data, "w", encoding="utf-8") as fd:
            json.dump(items_from_reader, fd, ensure_ascii=False, indent=2)
        logger.info("[LOCAL RAG] cache saved → %s", data)
    except Exception as e:
        logger.debug("cache save failed: %s", e)

def _file_version(path: str) -> str:
    """파일 변경을 반영하는 버전 식별자."""
    forced = os.getenv("LOCAL_RAG_FORCE_VERSION")
    if forced:
        logger.debug("Using forced version for %s: %s", path, forced)
        return str(forced)

    mode = (os.getenv("LOCAL_RAG_VERSION_MODE", "mtime") or "mtime").lower()
    try:
        if mode == "sha1":
            h = hashlib.sha1()
            with open(path, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    h.update(chunk)
            ver = h.hexdigest()[:12]
        else:
            st = os.stat(path)
            ver = f"{int(st.st_mtime)}_{st.st_size}"
        logger.debug("Computed version for %s: %s (mode=%s)", path, ver, mode)
        return ver
    except Exception as e:
        logger.warning("Version compute failed for %s: %s", path, e)
        return "na"


def _read_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


def _read_md(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _summarize_table_lines(lines: List[str], max_rows: int = 3) -> str:
    head = [ln.strip() for ln in lines[:max_rows] if ln.strip()]
    if not head:
        return "[표 요약]"
    sample = " / ".join(head)
    return f"[표 요약: {sample[:120]}...]"

_MD_FENCE_RE = re.compile(r"```.*?```", flags=re.DOTALL)
_MD_TABLE_LINE_RE = re.compile(r"^\s*\|.*\|\s*$")

def _split_markdown_to_chunks(
    text: str,
    min_chars: Optional[int] = None,
    max_chars: Optional[int] = None,
) -> List[str]:
    """
    규칙:
      - ```fence``` 코드블록은 '[코드 요약]'으로 치환
      - 파이프(|) 기반 테이블은 '[표 요약: ...]' 1줄로 치환
      - 빈줄 기준 문단 → 600~1,200자 범위로 병합/분할
    """
    min_chars = min_chars or _env_int("LOCAL_RAG_MIN_CHARS", 600)
    max_chars = max_chars or _env_int("LOCAL_RAG_MAX_CHARS", 1200)

    if os.getenv("LOCAL_RAG_CHUNK_MODE", "paragraph").lower() == "simple":
        # 롤백/토글 용: 매우 단순 라인 기반 (기존과 유사)
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        return lines

    # 1) 코드블록 요약 치환
    text = _MD_FENCE_RE.sub("[코드 요약]", text)

    # 2) 테이블 블록을 한 덩어리로 요약 치환
    lines = text.splitlines()
    out_lines, buf_table = [], []
    def _flush_table():
        nonlocal buf_table, out_lines
        if buf_table:
            out_lines.append(_summarize_table_lines(buf_table))
            buf_table = []
    for ln in lines:
        if _MD_TABLE_LINE_RE.match(ln):
            buf_table.append(ln)
        else:
            _flush_table()
            out_lines.append(ln)
    _flush_table()
    text = "\n".join(out_lines)

    # 3) 빈줄 기준 문단
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    # 4) 병합(짧은 문단 → 다음 문단과 붙임)
    merged, buf = [], ""
    for p in paras:
        if len(buf) + len(p) + 1 < min_chars:
            buf = (buf + " " + p).strip() if buf else p
        else:
            if buf:
                merged.append(buf)
            buf = p
    if buf:
        merged.append(buf)

    # 5) 최대 길이 초과는 단어 경계로 슬라이스
    chunks: List[str] = []
    for m in merged:
        cur = m.strip()
        while len(cur) > max_chars:
            cut = cur[:max_chars]
            cut_idx = cut.rfind(" ")
            if cut_idx < max_chars * 0.5:  # 단어 경계가 너무 앞이면 하드컷
                cut_idx = max_chars
            chunks.append(cur[:cut_idx].strip())
            cur = cur[cut_idx:].strip()
        if cur:
            chunks.append(cur)
    return chunks

def _truthy_env(name: str) -> bool:
    v = (os.getenv(name) or "").strip().lower()
    return v in ("1", "true", "yes", "on")


def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def _truncate(s: str, max_chars_env: str = "LOCAL_RAG_MAX_TEXT_CHARS") -> str:
    max_chars = int(os.getenv(max_chars_env, "200000"))
    if max_chars > 0 and len(s) > max_chars:
        return s[:max_chars]
    return s


def _read_html(path: str) -> str:
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    if BeautifulSoup is not None:
        try:
            return BeautifulSoup(raw, "lxml").get_text(" ", strip=True)
        except Exception:
            return BeautifulSoup(raw, "html.parser").get_text(" ", strip=True)
    txt = re.sub(r"<[^>]+>", " ", raw)
    return re.sub(r"\s+", " ", txt).strip()


def _read_docx(path: str) -> str:
    if not docx:
        raise RuntimeError("python-docx 미설치")
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs)

def _read_pptx_titles_and_bullets(path: str) -> List[dict]:
    """
    python-pptx 사용 가능 시: 슬라이드 단위로 제목/불릿만 추출하고,
    표/차트는 요약 플레이스홀더로 대체. (슬라이드 단위 리스트 반환)
    사용 불가 시: 빈 리스트 반환하여 상위에서 unstructured fallback.
    """
    if not _PptxPresentation:
        return []
    try:
        prs = _PptxPresentation(path)
    except Exception as e:
        logger.warning("python-pptx load failed for %s: %s", path, e)
        return []

    items: List[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        # 제목 shape
        try:
            title = getattr(slide.shapes.title, "text", "") or ""
            if title and len(title.strip()) >= 2:
                texts.append(title.strip())
        except Exception:
            pass

        # 본문/불릿/기타 텍스트
        for shape in slide.shapes:
            try:
                # 표
                if hasattr(shape, "has_table") and shape.has_table:
                    texts.append("[표 요약]")
                    continue
                # 차트
                if hasattr(shape, "chart"):
                    texts.append("[차트 요약]")
                    continue
                # 텍스트 프레임
                has_tf = bool(getattr(shape, "has_text_frame", False))
                tf = getattr(shape, "text_frame", None)
                if has_tf and tf is not None:
                    paragraphs = getattr(tf, "paragraphs", []) or []
                    for p in paragraphs:
                        t = "".join(r.text for r in getattr(p, "runs", []) or []) if getattr(p, "runs", None) else (getattr(p, "text", "") or "")
                        t = (t or "").strip()
                        if not t:
                            continue
                        if len(t) < 2:
                            continue
                        texts.append(t)
                else:
                    t = str(getattr(shape, "text", "") or "").strip()
                    if t and len(t) >= 2:
                        texts.append(t)
            except Exception:
                continue

        content = "\n".join(texts).strip()
        if content:
            items.append({
                "sheet": str(i),
                "row_index": i,
                "content": content,
            })
    return items


def _read_unstructured_elements(path: str) -> List[dict]:
    """
    PPTX, XLSX 등을 unstructured로 분해하여 요소(Element) 리스트로 반환.
    각 요소 dict 예시: {"sheet": "<슬라이드/시트/파트>", "row_index": <번호>, "content": "<텍스트>"}.
    """
    if not partition:
        raise RuntimeError("unstructured 라이브러리 미설치")
    try:
        elements = partition(filename=path)
    except Exception as e:
        logger.warning("Unstructured partition failed for %s: %s", path, e)
        return []

    items: List[dict] = []
    for i, element in enumerate(elements):
        # element.text가 없을 수 있음
        content = str(getattr(element, "text", "") or "").strip()
        if Table is not None and isinstance(element, Table):
            text_as_csv = getattr(element, "text_as_csv", "") or ""
            if text_as_csv:
                content = (content + "\n\n" + text_as_csv).strip()
        if not content:
            continue

        metadata = getattr(element, "metadata", None)
        part_name = getattr(metadata, "page_number", None) or getattr(metadata, "sheet_name", None) or "Part"

        items.append({
            "sheet": str(part_name),
            "row_index": i + 1,
            "content": content,
        })
    return items

def _pick_sample_indices(n: int, ratio: float) -> List[int]:
    import math
    k = max(1, int(math.ceil(n * max(0.01, min(1.0, ratio)))))
    # 항상 첫/끝 슬라이드는 포함
    base = {0, max(0, n - 1)}
    if n <= len(base):
        return sorted(base)
    # 간격 샘플링
    step = max(1, n // k)
    for i in range(0, n, step):
        base.add(i)
        if len(base) >= k:
            break
    return sorted(min(idx, n - 1) for idx in base)

def _pptx_extract_titles_bullets(path: str, sample_large: bool = False) -> List[dict]:
    """
    대용량 PPTX 전용: 제목+불릿만 추출, 표/차트 존재는 태그로 표기.
    sample_large=True면 슬라이드 20% 샘플링(기본), 항상 1/마지막 포함.
    반환은 items_from_reader 포맷: {"sheet": slide_no, "row_index": i, "content": "..."} 리스트
    """
    if _PptxPresentation is None:
        raise RuntimeError("python-pptx 미설치")

    prs = _PptxPresentation(path)
    slides = list(prs.slides)
    total = len(slides)
    ratio = _env_float("LOCAL_RAG_SAMPLE_RATIO", 0.2) if sample_large else 1.0
    pick = _pick_sample_indices(total, ratio)

    items: List[dict] = []
    for sidx in pick:
        slide = slides[sidx]
        texts: List[str] = []
        # 제목 placeholder
        title_shape = getattr(slide, "shapes", None)
        title = ""
        if title_shape is not None:
            for shp in slide.shapes:
                # 제목 후보: TitlePlaceholder or name에 'Title'
                name = str(getattr(shp, "name", "") or "")
                ph = getattr(shp, "placeholder_format", None)
                is_title = False
                if ph is not None:
                    tpe = getattr(ph, "type", None)
                    is_title = str(tpe).lower().find("title") >= 0
                if (not is_title) and ("title" in name.lower()):
                    is_title = True
                # 텍스트 프레임
                has_tf = bool(getattr(shp, "has_text_frame", False))
                tf = getattr(shp, "text_frame", None)
                if is_title and has_tf and tf is not None:
                    title = str(getattr(tf, "text", "") or "").strip()
                    break
        if title:
            texts.append(title)

        # 불릿 추출
        bullets: List[str] = []
        if title_shape is not None:
            for shp in slide.shapes:
                has_tf = bool(getattr(shp, "has_text_frame", False))
                tf = getattr(shp, "text_frame", None)
                if not (has_tf and tf):
                    continue
                paragraphs = getattr(tf, "paragraphs", []) or []
                for p in paragraphs:
                    # 제목 문단은 스킵
                    raw = str(getattr(p, "text", "") or "").strip()
                    if not raw:
                        continue
                    level = int(getattr(p, "level", 0) or 0)
                    if title and raw == title:
                        continue
                    # 불릿으로 간주: level>0 또는 문장 길이 2자 이상
                    if level > 0 or len(raw) >= 2:
                        bullets.append(raw)
        if bullets:
            texts.extend([f"- {b}" for b in bullets])

        # 표/차트 존재 플래그
        has_table = False
        has_chart = False
        num_shapes = 0
        if title_shape is not None:
            for shp in slide.shapes:
                num_shapes += 1
                if bool(getattr(shp, "has_table", False)):
                    has_table = True
                if getattr(shp, "chart", None) is not None:
                    has_chart = True
        if has_table:
            texts.append("[표 있음]")
        if has_chart:
            texts.append("[차트 있음]")

        content = "\n".join(t for t in texts if t).strip()
        if not content:
            continue

        items.append({
            "sheet": str(sidx + 1),
            "row_index": sidx + 1,
            "content": content,
        })

    logger.info("[LOCAL RAG] PPTX %s → slim extracted (slides=%d/%d)", Path(path).name, len(items), total)
    return items

def _read_pptx(path: str) -> List[dict]:
    """
    우선 python-pptx로 슬라이드 단위(제목+불릿) 추출 시도,
    실패/미설치 시 unstructured로 폴백.
    """
    if os.getenv("SKIP_PPTX", "0") in ("1", "true", "yes", "on"):
        logger.info("SKIP_PPTX=1 → _read_pptx skipped for %s", path)
        return []

    items = _read_pptx_titles_and_bullets(path)
    if items:
        return items
    return _read_unstructured_elements(path)



def _read_xlsx(path: str) -> List[dict]:
    return _read_unstructured_elements(path)


def _read_pdf_pages(path: str) -> List[str]:
    if not PdfReader:
        raise RuntimeError("PyPDF2 미설치")
    reader = PdfReader(path)
    max_pages = int(os.getenv("LOCAL_RAG_PDF_MAX_PAGES", "30"))
    pages: List[str] = []
    for i, p in enumerate(reader.pages, start=1):
        if i > max_pages:
            break
        try:
            txt = p.extract_text() or ""
            pages.append(txt)
        except Exception:
            pages.append("")
    return pages


def _file_uri(path: str) -> str:
    return Path(path).resolve().as_uri()  # file:/// 형태 보장

# ──────────────────────────────────────────────────────────────────────────────
# 변환: 로컬 파일 → web.json 아이템 배열
# ──────────────────────────────────────────────────────────────────────────────
def _to_webjson_items(path: str) -> List[dict]:
    ext = Path(path).suffix.lower()
    title = Path(path).name
    url_click = _file_uri(path)  # 사람이 눌러 열어볼 주소
    ver = _file_version(path)
    fetched_at = _now_iso()

    # 웹 파이프라인과 스키마 정렬: content_type 지정
    _ct_map = {
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".html": "text/html",
        ".htm": "text/html",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
    content_type = _ct_map.get(ext, "application/octet-stream")

    # 용량 가드 (+대용량 PPTX 샘플링/캐시)
    max_mb = _env_float("LOCAL_RAG_MAX_FILE_MB", 50.0)
    size_mb = None
    try:
        size_mb = os.path.getsize(path) / (1024 * 1024)
    except Exception:
        logger.debug("filesize check failed for %s", path)

    # 대용량 처리 조건
    is_large = bool(size_mb and max_mb > 0 and size_mb > max_mb)
    enable_sample = _truthy_env("LOCAL_RAG_SAMPLE_LARGE")  # 1/true이면 대용량 샘플링


     # 1) 파일 유형별 추출
    text: Optional[str] = None
    items_from_reader: Optional[List[dict]] = None
    try:
        if ext == ".pdf":
            if is_large:
                logger.info("[LOCAL RAG] skip large PDF (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            pages = _read_pdf_pages(path)
            items_from_reader = []
            for i, txt in enumerate(pages, start=1):
                if (txt or "").strip():
                    items_from_reader.append({
                        "page_num": i,
                        "content": txt,
                    })

        elif ext in (".txt", ".md", ".markdown", ".html", ".htm", ".docx"):
            if is_large:
                logger.info("[LOCAL RAG] skip large text-like (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            if ext == ".txt":
                text = _read_txt(path)
            elif ext in (".md", ".markdown"):
                text = _read_md(path)
            elif ext in (".html", ".htm"):
                text = _read_html(path)
            elif ext == ".docx":
                text = _read_docx(path)

        elif ext == ".pptx":
            if is_large and enable_sample:
                # 캐시 조회
                cached = _cache_load(path)
                if cached is not None:
                    items_from_reader = cached
                    logger.info("[LOCAL RAG] sampled/slim extracted (cache) → %s", Path(path).name)
                else:
                    # 샘플 슬림 추출
                    items_from_reader = _pptx_extract_titles_bullets(path, sample_large=True)
                    _cache_save(path, items_from_reader or [])
                    logger.info("[LOCAL RAG] sampled/slim extracted (fresh) → %s", Path(path).name)
            elif is_large and not enable_sample:
                logger.info("[LOCAL RAG] skip large file (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            else:
                # 일반 크기 PPTX는 기존 파이프라인(unstructured) 또는 슬림 추출 중 선택
                prefer_slim = _truthy_env("LOCAL_RAG_PPTX_SLIM")
                if prefer_slim and _PptxPresentation is not None:
                    items_from_reader = _pptx_extract_titles_bullets(path, sample_large=False)
                else:
                    items_from_reader = _read_unstructured_elements(path)

        elif ext == ".xlsx":
            if is_large:
                logger.info("[LOCAL RAG] skip large XLSX (> %.1f MB): %s (%.1f MB)", max_mb, path, size_mb or -1)
                return []
            items_from_reader = _read_unstructured_elements(path)

        else:
            logger.debug("[LOCAL RAG] unsupported extension skipped: %s", path)
            return []

    except RuntimeError as e:
        logger.warning("[LOCAL RAG] Reader dependency failed (RuntimeError): %s -> %s", path, e)
        return []
    except Exception as e:
        logger.warning("[LOCAL RAG] Extraction error: %s -> %s", path, e)


    # 2) 결과 통합 및 메타데이터 생성 (공통)
    final_items: List[dict] = []

    # 리스트 형태 결과(PDF, PPTX, XLSX)
    # 2) 결과 통합 및 메타데이터 생성 (공통)
    final_items: List[dict] = []

    # 리스트 형태 결과(PDF, PPTX, XLSX)
    if items_from_reader is not None and isinstance(items_from_reader, list):
        # 리스트 아이템별 content를 600~1,200자에 맞춰 재청킹
        minc = _env_int("LOCAL_RAG_MIN_CHARS", 600)
        maxc = _env_int("LOCAL_RAG_MAX_CHARS", 1200)

        for it in items_from_reader:
            part_label = str(it.pop("sheet", it.pop("page_num", "Part")))
            index_num = it.pop("row_index", it.pop("page_num", 1))
            content = (it.pop("content", "") or "").strip()
            if not content:
                continue

            # PPTX/XLSX는 그대로, PDF 등은 가볍게 문단 기준으로 쪼갬
            chunks = _split_markdown_to_chunks(content, minc, maxc)
            if not chunks:
                continue

            for j, ch in enumerate(chunks, start=1):
                ch = _truncate(ch, "LOCAL_RAG_MAX_TEXT_CHARS")
                final_items.append({
                    "title": f"{title} ({part_label}, Index: {index_num}, Chunk {j})",
                    "url": f"{url_click}#part={part_label}&index={index_num}",
                    "source": f"{url_click}__s_{part_label}__r_{index_num}__c_{j}__v_{ver}",
                    "content": ch,
                    "content_type": content_type,
                    "fetched_at": fetched_at,
                })


    # 단일 텍스트 결과(TXT, MD, HTML, DOCX)
    elif text is not None:
        content = _truncate(text or "", "LOCAL_RAG_MAX_TEXT_CHARS")
        if content.strip():
            minc = _env_int("LOCAL_RAG_MIN_CHARS", 600)
            maxc = _env_int("LOCAL_RAG_MAX_CHARS", 1200)
            # MD/HTML/TXT/DOCX → MD 규칙 기반으로 통일 처리
            chunks = _split_markdown_to_chunks(content, minc, maxc)
            for j, ch in enumerate(chunks, start=1):
                final_items.append({
                    "title": f"{title} (Chunk {j})",
                    "url": url_click,
                    "source": f"{url_click}__c_{j}__v_{ver}",
                    "content": ch,
                    "content_type": content_type,
                    "fetched_at": fetched_at,
                })

    if not final_items:
        logger.debug("[LOCAL RAG] File yielded no extractable content: %s", path)
        return []

    return final_items

# ──────────────────────────────────────────────────────────────────────────────
# 엔트리: globs → web.json 생성
# ──────────────────────────────────────────────────────────────────────────────
def build_webjson_from_local(globs: List[str], out_dir: str) -> str:
    logger.info("[LOCAL RAG] CWD: %s", os.getcwd())
    logger.info("[LOCAL RAG] Received globs: %s", globs)

    Path(out_dir).mkdir(parents=True, exist_ok=True)
    files: List[str] = []

    for g in globs or []:
        g = os.path.expandvars(os.path.expanduser(g))
        logger.debug("[LOCAL RAG] Expanded glob pattern: %s", g)
        matched = glob.glob(g, recursive=True)
        files.extend(matched)
        logger.info("[LOCAL RAG] Pattern %s matched %d files.", g, len(matched))

    files = sorted({f for f in files if os.path.isfile(f)})
    logger.info("[LOCAL RAG] Total unique files found: %d", len(files))

    items: List[dict] = []
    for f in files:
        try:
            items.extend(_to_webjson_items(f))
        except Exception as e:
            logger.warning("[LOCAL RAG] local ingest 실패: %s -> %s", f, e)

    # 빈 content 제거
    items = [it for it in items if (it.get("content") or "").strip()]

    # 중복 source 제거(최초 등장 우선)
    dedup_seen, dedup_items = set(), []
    for it in items:
        sid = it.get("source")
        if not sid or sid in dedup_seen:
            continue
        dedup_seen.add(sid)
        dedup_items.append(it)
    items = dedup_items

    # 디버그 샘플
    uniq_sources = {it.get("source") for it in items}
    logger.info("[LOCAL RAG] files=%d items=%d unique_sources=%d", len(files), len(items), len(uniq_sources))
    if items:
        def _pretty_src(src: str) -> str:
            s = unquote(src or "")
            if "__v_" in s:
                s = s.split("__v_")[0]
            if "__p_" in s:
                s = s.split("__p_")[0]
            return s

        sample = items[:3]
        sample_titles = [it.get("title", "") for it in sample]
        sample_sources = [_pretty_src(it.get("source", "")) for it in sample]
        sample_urls = [unquote(it.get("url", "")) for it in sample]

        logger.debug("[LOCAL RAG] sample titles : %s", sample_titles)
        logger.debug("[LOCAL RAG] sample sources: %s", sample_sources)
        logger.debug("[LOCAL RAG] sample urls   : %s", sample_urls)

    ts = datetime.now().strftime("%Y_%m_%d_%H%M%S")
    out_path = os.path.join(out_dir, f"local_{ts}.json")
    with open(out_path, "w", encoding="utf-8") as fp:
        json.dump(items, fp, ensure_ascii=False, indent=2)
    logger.info("[LOCAL RAG] web.json saved → %s", out_path)
    return out_path

# ──────────────────────────────────────────────────────────────────────────────
# 파이프: web.json → Chroma 적재 + 미리보기
# ──────────────────────────────────────────────────────────────────────────────
def ingest_local_files(
    globs: List[str],
    namespace: str,
    persist_directory: str | None,
    topic_slug: str,
    root_dir: str,
    add_web_pages_json_to_chroma=None,
    web_page_json_to_documents=None,
) -> Tuple[List[str], List[Document], int]:
    """
    반환: (생성 JSON 경로들, 미리보기 Documents, 인덱싱된 청크 수 합)
    """
    # 연구 요약(findings) 포함 옵션
    if _truthy_env("INCLUDE_FINDINGS_IN_VECTOR"):
        slug = (topic_slug or os.getenv("TOPIC_SLUG") or "default").strip()
        findings_pattern = os.path.join(root_dir, "research", slug, "round-*-findings.md")
        globs = list(globs or [])
        globs.append(findings_pattern)
        logger.info("[LOCAL RAG] findings included → %s", findings_pattern)

    if not globs:
        logger.info("[LOCAL RAG] no globs provided → skip ingest")
        return ([], [], 0)

    res_dir = os.path.join(root_dir, "resources", topic_slug or "default")
    json_path = build_webjson_from_local(globs, res_dir)

    chunk_total = 0
    if add_web_pages_json_to_chroma is not None:
        try:
            _orig, chunk_count = add_web_pages_json_to_chroma(
                json_path, namespace=namespace, persist_directory=persist_directory
            )
            chunk_total += int(chunk_count or 0)
            logger.info("[LOCAL RAG] added to chroma: chunks=%s (ns=%s, dir=%s)", chunk_count, namespace, persist_directory)
        except Exception as e:
            logger.warning("[LOCAL RAG] add_web_pages_json_to_chroma(local) 실패: %s", e)

    docs_preview: List[Document] = []
    if web_page_json_to_documents is not None:
        try:
            docs_preview = web_page_json_to_documents(json_path)[:8]
            logger.debug("[LOCAL RAG] preview docs: %d", len(docs_preview))
        except Exception as e:
            logger.warning("[LOCAL RAG] preview build(local) 실패: %s", e)

    return ([json_path], docs_preview, chunk_total)
