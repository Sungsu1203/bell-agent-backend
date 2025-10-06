# tools/local_rag.py
from __future__ import annotations
import os, re, json, glob, hashlib
from pathlib import Path
from datetime import datetime
from langchain_core.documents import Document

from typing import Any, List, Tuple

from urllib.parse import unquote

# 선택 의존성: 있으면 사용, 없으면 경고만
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import docx  # python-docx
except Exception:
    docx = None

# try:
#     from pptx import Presentation  # python-pptx
# except Exception:
#     Presentation = None

try:
    from bs4 import BeautifulSoup  # beautifulsoup4
except Exception:
    BeautifulSoup = None

def _file_version(path: str) -> str:
    """
    파일 변경을 반영하는 버전 식별자.
    - 기본: mtime + size
    - ENV LOCAL_RAG_VERSION_MODE=sha1 이면 내용 해시 사용
    - ENV LOCAL_RAG_FORCE_VERSION 이 설정되면 그 값을 그대로 사용(강제 리인덱싱용)
    """
    forced = os.getenv("LOCAL_RAG_FORCE_VERSION")
    if forced:
        return str(forced)

    mode = os.getenv("LOCAL_RAG_VERSION_MODE", "mtime").lower()
    try:
        if mode == "sha1":
            h = hashlib.sha1()
            with open(path, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    h.update(chunk)
            return h.hexdigest()[:12]
        else:
            st = os.stat(path)
            return f"{int(st.st_mtime)}_{st.st_size}"
    except Exception:
        return "na"


def _read_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _read_md(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _read_html(path: str) -> str:
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    if BeautifulSoup:
        try:
            return BeautifulSoup(raw, "lxml").get_text(" ", strip=True)
        except Exception:
            return BeautifulSoup(raw, "html.parser").get_text(" ", strip=True)
    txt = re.sub(r"<[^>]+>", " ", raw)
    return re.sub(r"\s+", " ", txt).strip()

def _read_docx(path: str) -> str:
    if not docx:
        raise RuntimeError("python-docx 미설치")
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs)


# 그룹 도형 판별용
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    MSO_SHAPE_TYPE = None  # 라이브러리 미설치 시에도 코드가 동작하도록

def _shape_texts(shape: Any) -> List[str]:
    texts: List[str] = []

# 텍스트 프레임 → 표 셀 → 차트 제목 → 대체텍스트 → 그룹 재귀 순서로, 도형 하나에서 나올 수 있는 텍스트를 빠짐없이 수집

    # 1) 일반 도형의 텍스트 프레임
    tf = getattr(shape, "text_frame", None)
    if tf is not None:
        texts.append("\n".join((p.text or "") for p in tf.paragraphs))

    # 2) 표(table) 셀 텍스트
    tbl = getattr(shape, "table", None)
    if tbl is not None:
        for row in tbl.rows:
            for cell in row.cells:
                texts.append(cell.text or "")

    # 3) 차트 제목
    chart = getattr(shape, "chart", None)
    if chart is not None and getattr(chart, "has_title", False):
        tft = getattr(chart.chart_title, "text_frame", None)
        if tft is not None:
            texts.append("\n".join((p.text or "") for p in tft.paragraphs))

    # 3-b) (옵션) 그림 대체텍스트(alt text)
    alt = getattr(shape, "alternative_text", None)
    if isinstance(alt, str) and alt.strip():
        texts.append(alt)

    # 4) 그룹 도형 내부 재귀
    shape_type = getattr(shape, "shape_type", None)
    if (MSO_SHAPE_TYPE and shape_type == MSO_SHAPE_TYPE.GROUP) or hasattr(shape, "shapes"):
        for inner in getattr(shape, "shapes", []):
            texts.extend(_shape_texts(inner))

    return texts

def _read_pptx(path: str) -> str:
    # 함수 내부 임포트로 OptionalCall 경고 방지
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx 미설치") from e

    prs = Presentation(path)

    parts: List[str] = []
    for slide in prs.slides:
        # 슬라이드의 모든 도형에서 텍스트 수집
        for shape in slide.shapes:
            parts.extend(_shape_texts(shape))

        # (옵션) 슬라이드 노트 텍스트도 수집
        notes = getattr(slide, "notes_slide", None)
        if notes is not None:
            ntf = getattr(notes, "notes_text_frame", None)
            if ntf is not None:
                parts.append(ntf.text or "")

    # 공백 정리 + 빈 문자열 제거
    cleaned: List[str] = []
    for t in parts:
        if not t:
            continue
        s = re.sub(r"\s+", " ", t).strip()
        if s:
            cleaned.append(s)

    return "\n".join(cleaned)


def _read_pdf_pages(path: str) -> List[str]:
    if not PdfReader:
        raise RuntimeError("pypdf 미설치")
    reader = PdfReader(path)

    max_pages = int(os.getenv("LOCAL_RAG_PDF_MAX_PAGES", "30"))
    pages: List[str] = []
    for i, p in enumerate(reader.pages, start=1):
        if i > max_pages:
            break
        try:
            pages.append(p.extract_text() or "")
        except Exception:
            pages.append("")
    return pages


# ── 추가: 해시 & 파일 URI ───────────────────────────────────────────────
# def _sha1(text: str) -> str:
#     return hashlib.sha1((text or "").encode("utf-8", "ignore")).hexdigest()[:12]

def _file_uri(path: str) -> str:
    return Path(path).resolve().as_uri()  # file:/// 형태 보장
# ────────────────────────────────────────────────────────────────────────

def _to_webjson_items(path: str) -> List[dict]:
    ext   = Path(path).suffix.lower()
    title = Path(path).name
    # 클릭용 URL (사람이 눌러 열어볼 주소)
    url_click = _file_uri(path)
    ver = _file_version(path)

    # 용량 가드
    max_mb = float(os.getenv("LOCAL_RAG_MAX_FILE_MB", "50"))
    try:
        if max_mb > 0:
            size_mb = os.path.getsize(path) / (1024 * 1024)
            if size_mb > max_mb:
                print(f"[LOCAL RAG] skip large file (> {max_mb}MB): {path} ({size_mb:.1f}MB)")
                return []
    except Exception:
        pass

    if ext == ".pdf":
        pages = _read_pdf_pages(path)
        items = []
        for i, txt in enumerate(pages, start=1):
            if not (txt or "").strip():
                continue
            # 디듀프 키(Chroma 중복 판정용)는 접미사 버전 사용
            source_key = f"{url_click}__p_{i}__v_{ver}"   # ❗쿼리/프래그먼트 없이 접미사
            # 클릭용 URL은 그대로(원하면 보기 편하게 #page 유지)
            url_for_click = f"{url_click}#page={i}"
            max_chars = int(os.getenv("LOCAL_RAG_MAX_TEXT_CHARS", "200000"))
            ct = txt or ""
            if len(ct) > max_chars:
                ct = ct[:max_chars]

            items.append({
                "title": f"{title} (p.{i})",
                "url": url_for_click,
                "source": source_key,
                "content": ct
            })
        if not items:
            items = [{"title": title, "url": url_click, "source": f"{url_click}__v_{ver}", "content": ""}]
        return items

    # PDF 외 파일
    if ext in [".txt"]:
        text = _read_txt(path)
    elif ext in [".md", ".markdown"]:
        text = _read_md(path)
    elif ext in [".html", ".htm"]:
        text = _read_html(path)
    elif ext in [".docx"]:
        text = _read_docx(path)
    elif ext in [".pptx"]:
        text = _read_pptx(path)
    else:
        return []

    max_chars = int(os.getenv("LOCAL_RAG_MAX_TEXT_CHARS", "200000"))
    content = (text or "")
    if len(content) > max_chars:
        content = content[:max_chars]

    return [{
        "title": title,
        "url": url_click,                 # 클릭용
        "source": f"{url_click}__v_{ver}",# 디듀프 키(접미사 버전)
        "content": content
    }]

def build_webjson_from_local(globs: List[str], out_dir: str) -> str:
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    files: List[str] = []
    for g in globs:
        g = os.path.expandvars(os.path.expanduser(g))  # ← 추가
        files.extend(glob.glob(g, recursive=True))
    files = sorted({f for f in files if os.path.isfile(f)})

    items: List[dict] = []
    for f in files:
        try:
            items.extend(_to_webjson_items(f))
        except Exception as e:
            print(f"[WARN] local ingest 실패: {f} -> {e}")

    # 빈 content 제거
    items = [it for it in items if (it.get("content") or "").strip()]

    # 디버그 요약 (web_rag와 조응되는 JSON 스키마 기준)
    uniq_sources = {it.get("source") for it in items}
    print(f"[LOCAL RAG] files={len(files)} items={len(items)} unique_sources={len(uniq_sources)}")
    if items:
        # 사람이 읽기 편하도록 퍼센트 인코딩 해제 + 접미사 제거
        def _pretty_src(src: str) -> str:
            s = unquote(src or "")
            # 페이지/버전 접미사는 로그에서만 제거 (저장값은 그대로)
            if "__v_" in s:
                s = s.split("__v_")[0]
            if "__p_" in s:
                s = s.split("__p_")[0]
            return s

        sample = items[:3]
        sample_titles  = [it.get("title", "") for it in sample]
        sample_sources = [_pretty_src(it.get("source", "")) for it in sample]
        sample_urls    = [unquote(it.get("url", "")) for it in sample]

        print(f"[LOCAL RAG] sample titles : {sample_titles}")
        print(f"[LOCAL RAG] sample sources: {sample_sources}")
        print(f"[LOCAL RAG] sample urls   : {sample_urls}")

    ts = datetime.now().strftime("%Y_%m_%d_%H%M%S")
    out_path = os.path.join(out_dir, f"local_{ts}.json")
    with open(out_path, "w", encoding="utf-8") as fp:
        json.dump(items, fp, ensure_ascii=False, indent=2)
    return out_path

def ingest_local_files(
    globs: List[str],
    namespace: str,
    persist_directory: str | None,
    topic_slug: str,
    root_dir: str,
    add_web_pages_json_to_chroma=None,
    web_page_json_to_documents=None,
) -> Tuple[List[str], List[Document], int]:
    """
    반환: (생성 JSON 경로들, 미리보기 Documents, 인덱싱된 청크 수 합)
    """
    if not globs:
        return ([], [], 0)

    res_dir = os.path.join(root_dir, "resources", topic_slug or "default")
    json_path = build_webjson_from_local(globs, res_dir)

    chunk_total = 0
    if add_web_pages_json_to_chroma:
        try:
            _orig, chunk_count = add_web_pages_json_to_chroma(
                json_path, namespace=namespace, persist_directory=persist_directory
            )
            chunk_total += int(chunk_count or 0)
        except Exception as e:
            print(f"[WARN] add_web_pages_json_to_chroma(local) 실패: {e}")

    docs_preview: List[Document] = []
    if web_page_json_to_documents:
        try:
            docs_preview = web_page_json_to_documents(json_path)[:8]
        except Exception as e:
            print(f"[WARN] preview build(local) 실패: {e}")

    return ([json_path], docs_preview, chunk_total)
